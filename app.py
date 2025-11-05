
import streamlit as st
from langgraph.graph import StateGraph, END
from typing import TypedDict, Optional, List
from langchain_core.runnables import Runnable
from serpapi import GoogleSearch
from vanna.remote import VannaDefault
from docx import Document
import tempfile
import os
from dotenv import load_dotenv
import json
import re
from openai import OpenAI
import pandas as pd
import sqlite3
from typing import Optional
import matplotlib.pyplot as plt
import networkx as nx
from io import BytesIO
from datetime import datetime
import json
from datetime import date, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import OpenAIEmbeddings
import tempfile
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.text.paragraph import Paragraph
from docx.table import Table
from docx.table import Table as DocxTable
import tempfile
import uuid
from docx.table import Table as DocxTable
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
import numpy as np
from pathlib import Path
from io import BytesIO
import base64
from time import sleep
from urllib.parse import urlparse
import io
import re
from difflib import get_close_matches

st.set_page_config(layout="wide")

load_dotenv()

#os.environ["LANGCHAIN_TRACING_V2"]="true"
#os.environ["LANGCHAIN_API_KEY"]=os.getenv("LANGCHAIN_API_KEY")

embedding = OpenAIEmbeddings(openai_api_key=os.getenv("OPENAI_API_KEY"))
faiss_index = FAISS.load_local("faiss_index", embedding, allow_dangerous_deserialization=True)

# Keywords that usually indicate monetary columns
money_keywords = ["loss", "premium", "amount", "cost", "ibnr", "ult", "total", "claim", "reserve", "payment"]

# ---- Vanna Setup ----
vanna_api_key = os.getenv("vanna_api_key")
vanna_model_name = os.getenv("vanna_model_name")
vn_model = VannaDefault(model=vanna_model_name, api_key=vanna_api_key)
vn_model.connect_to_sqlite('Actuarial_Data.db')

# ---- Config ----
#openai.api_key = os.getenv("TOGETHER_API_KEY")
#openai.base_url = "https://api.together.xyz"
serpapi_key = os.getenv("SERPAPI_API_KEY")

# ---Open AI LLM---

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

def call_llm(prompt: str) -> str:
    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an intelligent AI assistant."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.1,
            max_tokens=500
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"OpenAI call failed: {e}"


def last(old, new):
    return new

# ---- Define LangGraph State ----
class GraphState(TypedDict):
    user_prompt: str
    doc_loaded: bool
    document_path: Optional[str]

    # Prompts
    vanna_prompt: Optional[str]
    fuzzy_prompt: Optional[str]

    # Routing
    route: Optional[str]

    # SQL-related
    sql_result: Optional[pd.DataFrame]
    sql_query: Optional[str]

    # Document-related
    updated_doc_path: Optional[str]
    header_candidate: Optional[str]          # fuzzy-matched header
    table_candidate_index: Optional[int]     # fuzzy-matched table index
    header_updated: Optional[str]            # final confirmed header
    table_index_updated: Optional[int]       # final confirmed table index
    #candidate_tables: Optional[list[dict]]  # store top-n candidates
    updated_doc_key: Optional[str]           # unique Streamlit key for download button
    preview_df: Optional[list[dict]]         # serializable preview rows (list of dicts): table from doc
    preview_df_columns: Optional[list[str]]  # column names for preview_df

    # External search
    web_links: Optional[List[str]]

    # Visualization / summaries
    chart_info: Optional[dict]
    comparison_summary: Optional[str]
    general_summary: Optional[str]

    # FAISS Knowledge base
    faiss_summary: Optional[str]
    faiss_sources: Optional[list[tuple[str, str]]]
    faiss_images: Optional[list[dict]]

    # Reconciliation Agent
    # Reconciliation / upload fields (add to GraphState)
    uploaded_file1: Optional[object] = None       # the UploadedFile object (or None)
    uploaded_file2: Optional[object] = None       # the UploadedFile object (or None)
    uploaded_file1_path: Optional[str] = None     # temp file path on disk (or None)
    uploaded_file2_path: Optional[str] = None     # temp file path on disk (or None)
    uploaded_file1_is_excel: bool = False         # True if file1 appears excel/csv-like
    uploaded_file2_is_excel: bool = False         # True if file2 appears excel/csv-like
    reconcile_df: Optional[pd.DataFrame]
    missing_in_A: Optional[pd.DataFrame]
    missing_in_B: Optional[pd.DataFrame]
    variance_summary: Optional[dict]
    variance_commentary: Optional[str]
    comp_type: Optional[str]


def get_schema_description(db_path: str) -> str:
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()

    schema_str = ""
    cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
    tables = cursor.fetchall()

    for table_name, in tables:
        cursor.execute(f"PRAGMA table_info({table_name});")
        cols = cursor.fetchall()
        col_names = [col[1] for col in cols]
        schema_str += f"\n- {table_name}: columns = {', '.join(col_names)}"

    conn.close()
    return schema_str.strip()

def load_qs_pairs():
    with open("vanna_advanced_sql_pairs.txt", "r") as f:
        text = f.read()
    pairs = re.findall(r'question="(.*?)",\s*sql="""(.*?)"""', text, re.DOTALL)
    return [{"question": q.strip(), "sql": s.strip()} for q, s in pairs]

QSPairs = load_qs_pairs()
qs_examples = "\n".join(
    f"Q: {pair['question']}\nSQL: {pair['sql']}" for pair in QSPairs[:7]  # Limit to 7 to avoid token overflow
)


documentation = """
PnC_Data Table:
- Reserve Class contains insurance business lines such as 'Property', 'Casualty', 'Marine', 'Motor', etc.
- Exposure Year refers to the year in which the insured risk was exposed to potential loss.
- RI Type identifies whether the record is 'Gross' or one of the reinsurance types such as 'Ceded - XOL', 'Ceded - QS', 'Ceded - CAP', 'Ceded - FAC', or 'Ceded - Others'.
- Branch indicates the geographical business unit handling the contract, e.g., 'Europe', 'LATAM', 'North America'.
- Loss Type captures the nature of the loss, and may be one of: 'ATT', 'CAT', 'LARGE', 'THREAT', or 'Disc'.
- Underwriting Year represents the year in which the policy was underwritten or originated.
- Incurred Loss represents the total loss incurred to date, including paid and case reserves.
- Paid Loss is the portion of the Incurred Loss that has already been settled and paid out.
- IBNR is calculated as the difference between Ultimate Loss and Incurred Loss.
- Ultimate Loss is the projected final value of loss.
- Ultimate Premium refers to the projected premium expected to be earned.
- Loss Ratio is calculated as Ultimate Loss divided by Ultimate Premium.
- AvE Incurred = Expected - Actual Incurred.
- AvE Paid = Expected - Actual Paid.
- Budget Premium is the forecasted premium for budgeting.
- Budget Loss is the projected loss for budgeting.
- Earned Premium is the portion of the premium that has been earned.
- Case Reserves = Incurred Loss - Paid Loss.
"""

STATE_KEYS_SET_AT_ENTRY = [
#     "user_prompt", 
#     "doc_loaded", 
#     "document_path", 

#     # Prompts
#     "vanna_prompt", 
#     "fuzzy_prompt",

#     # Routing
#     "route",

#     # SQL-related
#     "sql_result",
#     "sql_query",

#     # Document-related
#     "updated_doc_path",
#     "header_candidate",
#     "table_candidate_index",
#     "header_updated",
#     "table_index_updated",
#     "candidate_tables",
#     "updated_doc_key",
#     "preview_df",   
#     "preview_df_columns",

#     # External search
#     "web_links",

#     # Visualization / summaries
#     "chart_info",
#     "comparison_summary",
#     "general_summary",

#     # FAISS Knowledge base
#     "faiss_summary", 
#     "faiss_sources",
#     "faiss_images"
    ## Reconciliation / upload fields (add to GraphState)
    # uploaded_file1: Optional[object] = None       # the UploadedFile object (or None)
    # uploaded_file2: Optional[object] = None       # the UploadedFile object (or None)
    # uploaded_file1_path: Optional[str] = None     # temp file path on disk (or None)
    # uploaded_file2_path: Optional[str] = None     # temp file path on disk (or None)
    # uploaded_file1_is_excel: bool = False         # True if file1 appears excel/csv-like
    # uploaded_file2_is_excel: bool = False         # True if file2 appears excel/csv-like
    # "reconcile_df",
    # "missing_in_A",
    # "missing_in_B",
    # "variance_summary",
    # "variance_commentary",
    # "reconcile_source_files",
    # "comp_type"
]

# --- short-circuit: if both uploads exist and look like excel/csv, route deterministically to comp ---
def _is_excel_like(uploaded_file):
    try:
        if not uploaded_file:
            return False
        name = getattr(uploaded_file, "name", "") or str(uploaded_file)
        name = name.lower()
        return name.endswith((".xls", ".xlsx", ".csv"))
    except Exception:
        return False


def prune_state(state: GraphState, exclude: List[str]) -> dict:
    return {k: v for k, v in state.items() if k not in exclude}

# ---- Router Node (with prompt generation) ----
class RouterNode(Runnable):
    def invoke(self, state: GraphState, config=None) -> GraphState:
        #doc_flag = "yes" if state['doc_loaded'] else "no"
        doc_flag1 = "yes" if state.get(uploaded_file1_is_excel) else "no"
        doc_flag2 = "yes" if state.get(uploaded_file2_is_excel) else "no"
        doc1_exist = "yes" if state.get(file1_path) else "no"
        doc2_exist = "yes" if state.get(file2_path) else "no"
        schema = get_schema_description('Actuarial_Data.db')

        router_prompt = f"""
        You are an intelligent routing agent. Your job is to:
        1. Choose one of the paths: "sql", "search", "document", "comp", "faissdb" based on the user prompt.
        2. Choose:
        - "sql" if the user is asking a question about structured insurance data (e.g. claims, premiums, reserves, IBNR, trends, comparisons across years or products) or something that can be answered from the following database schema:
            {schema}
        - Use this additional documentation to better understand column meanings:
          {documentation}
        - Additionally, here are some examples of SQL-style questions and their corresponding queries (QSPairs):
          {qs_examples}
        -EVEN IF the user also says things like "plot", "draw", "visualize", "graph", "bar chart", etc. ‚Äî that still means they want structured data **along with** a chart. SO route it to SQL
            Example: "Show me IBNR over years and plot a bar chart" ‚Üí route = "sql"
        -Route it to "sql" if queries includes the below mentioned:
            - Asks for trends, breakdowns, or aggregations of internal metrics (e.g., IBNR, reserves, severity, premiums, earned/ultimate loss)
            - Ask for trends **within internal data only**
            - Compares **internal data over time or segments** (e.g., years, lines of business, regions)
            - Ask for charts or visualizations ("plot", "bar chart", etc.)
            - Does NOT involve external benchmarking
            Even if the prompt includes words like "compare" or "change", still route to SQL if the context is strictly internal.
        -If the route is "sql", include vanna_prompt, but don't include fuzzy_prompt
            -(eg: User Prompt is "Show me exposure year wise incurred loss and plot a graph", then 
            -vanna_prompt will be "Shoe me exposure year wise incurred loss".
            -Your work is to remove the noise and focus only on things that are required to generate sql query from vanna. SO remove all the extra stuffs out of the user prompt.


        3. "document" ONLY if one document is uploaded AND the question involves updating/reading a document.
        -If the route is "document", DO NOT include vanna_prompt or fuzzy_prompt.
        -Check the status of the flag, only doc1_exist should be "yes" while doc2_exist should be "no". Status for
            -Doc_flag1 = {doc1_exist}
            -Doc_flag2 = {doc2_exist}
        
        4. Choose "search" if:
            - The user is asking about general or external information
            - Involves real-time info, news, global economic trends, regulations
            - The query cannot be answered by internal structured data or uploaded document
        - If the route is "search", DO NOT include vanna_prompt or fuzzy_prompt.


        5. Choose "comp" if both Doc_flag1 and Doc_flag2 are attached and both are excel files. Status for
        -Doc_flag1 = {doc_flag1}
        -Doc_flag2 = {doc_flag2}
        - Return Vanna_prompt: same as user_prompt
        
        6.Choose "comp" when the user is comparing internal data against external data, competitors, or industry benchmarks.
            Examples include peer review, benchmarking, market positioning, or competitive ratios.

            Trigger words/phrases (especially relevant for Actuarial & Finance users):
            - "industry benchmark"
            - "market average"
            - "how do we compare to..."
            - "peer comparison"
            - "market trend vs ours"
            - "against competitors"
            - "vs industry"
            - "benchmarking analysis"
            - "loss ratio gap with peers"
            - "pricing differential with market"
            - "expense ratio compared to competition"
            - "where do we stand in market"
            - "relative to industry"
            - "competitive advantage in reserves"
            - "our severity vs others"
            - "compare to S&P average" / "AM Best stats" / "regulatory benchmark"
        -(e.g.,
         1. If User_Prompt is "Compare IBNR trends with industry benchmarks for exposure year 2025"
            - Return Vanna_prompt: "Show IBNR trends for exposure year 2025"
         2. If User_Prompt is "what are the incurred loss trends as compared to axaxl competitors"
            - Return Vanna_prompt: "what are the incurred loss trends"
         3. If User_Prompt is "what are the expected loss trends as compared to market average"
            - Return Vanna_prompt: "Swhat are the ultimate loss trends"

        -Do not include fuzzy_prompt
        -Only include relevant columns in vanna_prompt. Do not include ClaimNumber or ID columns unless the user specifically asks for them.

       
        7. Choose "faissdb" when:
        - The prompt asks about the Sparta platform, Earmark Template, Branch Adjustment Template/Module, Projects in Sparta, or any internal process or documentation.
        - The user seems to be referring to internal workflows, or knowledge base content.
        -Example prompts that should be routed to `"faissdb"`:
            - "What are the steps in the Branch Adjustment Module?"
            - "Explain how Earmark Template is used in our process."
            - "Can you summarize Projects in Sparta?"


        Return output strictly in valid JSON format using double quotes and commas properly.
        DO NOT include any trailing commas. Your JSON must be parseable by Python's json.loads().

        Examples:

        For SQL:
        {{
            "route": "sql",
            "vanna_prompt": "Show IBNR trends for exposure year 2025"
        }}

        For Document:
        {{
            "route": "document",
        }}

        For Comp:
        {{
             "route": "comp",
             "vanna_prompt": "Show IBNR trends for exposure year 2025"
        }}

        For Search:
        {{
            "route": "search"
        }}

        For faissdb 
        {{
        "route": "faissdb"
        }}

        User Prompt: "{state['user_prompt']}"
        """
        #Document Uploaded: {doc_flag}
        
        if _is_excel_like(state.get("uploaded_file1")) and _is_excel_like(state.get("uploaded_file2")):
            # Deterministic routing: both attachments present and are excel-like
            parsed = {
                "route": "comp",
                "vanna_prompt": state.get("user_prompt"),
                "fuzzy_prompt": None
            }
            chart_info = None

        else:
            try:
                response = call_llm(router_prompt)
                #st.write("Route:", response)

                match = re.search(r'{.*}', response, re.DOTALL)
                if match:
                    parsed = json.loads(match.group())
                    chart_info = parsed.get("chart_info")
                else:
                    st.warning("LLM did not return valid JSON. Routing to 'search' as fallback.")
                    parsed = {"route": "search"}

            except Exception as e:
                st.error(f"[RouterNode] LLM call failed: {e}")
                parsed = {"route": "search"}

            # ‚úÖ Enforce safety: remove vanna_prompt
            if parsed.get("route") == "document":
                parsed["fuzzy_prompt"] = state["user_prompt"]   # alias
                parsed["vanna_prompt"] = None                   # will be set later
            elif parsed.get("route") not in ["sql", "comp", "faissdb"]:
                parsed["vanna_prompt"] = None
                parsed["fuzzy_prompt"] = None

            # Ensure chart_info is only kept for SQL route
            if parsed.get("route") != "sql":
                chart_info = None
        
        return {
            **prune_state(state, STATE_KEYS_SET_AT_ENTRY),
            "route": parsed.get("route"),
            "vanna_prompt": parsed.get("vanna_prompt"),
            "fuzzy_prompt": parsed.get("fuzzy_prompt"),
            "chart_info": chart_info,

            # Reset doc-specific keys at routing
            "header_candidate": None,
            "table_candidate_index": None,
            "header_updated": None,
            "table_index_updated": None,

            #Reconciliation Agent
            "reconcile_df": None,
            "missing_in_A": None,
            "missing_in_B": None,
            "variance_summary": None,
            "variance_commentary": None,
            "reconcile_source_files": None,
        }
    
# ---- Vanna SQL Node ----

def get_user_chart_type(prompt: str) -> Optional[str]:
    prompt = prompt.lower()
    if "bar chart" in prompt or "bar graph" in prompt:
        return "bar"
    elif "line chart" in prompt or "line graph" in prompt:
        return "line"
    elif "pie chart" in prompt or "pie graph" in prompt:
        return "pie"
    return None


def suggest_chart(df: pd.DataFrame) -> Optional[dict]:
    sample_data = df.head(5).to_dict(orient="list")
    prompt = f"""
    You are a data visualization assistant.

    Here is the top of a pandas DataFrame:
    {json.dumps(sample_data, indent=2)}

    Your task:
    - Identify a good chart (bar, line, or pie) that best represents this data.
    - Choose 1 column for the x-axis (categorical or time-based), and 1 or more numeric columns for the y-axis.
    - If multiple y columns are appropriate (e.g. IBNR, IncurredLoss), return them as a list.

    Return your answer in JSON like:
    {{ "type": "bar", "x": "ExposureYear", "y": ["IncurredLoss", "IBNR"] }}

    If no chart is suitable, return: "none"
    """

    reply = call_llm(prompt)
    match = re.search(r'{.*}', reply, re.DOTALL)
    if match:
        try:
            return json.loads(match.group())
        except:
            return None
    return None


def plot_chart(df: pd.DataFrame, chart_info: dict):
    chart_type = chart_info.get("type", "bar")
    x = chart_info.get("x")
    y = chart_info.get("y")

    if isinstance(y, str):
        y = [y]  # Make it a list

    df_columns = list(df.columns)
    def match_col(col_name):
        for c in df_columns:
            if col_name.lower().replace(" ", "") in c.lower().replace(" ", ""):
                return c
        return None

    x_col = match_col(x)
    y_cols = [match_col(col) for col in y if match_col(col)]

    if not x_col or not y_cols:
        st.warning(f"Invalid chart columns: {x}, {y}")
        return

    st.subheader(f"{chart_type.capitalize()} Chart: {', '.join(y)} vs {x}")

    if chart_type == "bar":
        st.bar_chart(df.set_index(x_col)[y_cols])
    elif chart_type == "line":
        st.line_chart(df.set_index(x_col)[y_cols])
    elif chart_type == "pie" and len(y_cols) == 1:
        fig, ax = plt.subplots()
        df.groupby(x_col)[y_cols[0]].sum().plot.pie(ax=ax, autopct='%1.1f%%')
        ax.set_ylabel('')
        st.pyplot(fig)
    else:
        st.warning("Pie chart supports only one y column.")

def vanna_node(state: GraphState) -> GraphState:
    # Use user_prompt if vanna_prompt is not available

    schema_desc = get_schema_description('Actuarial_Data.db')
    raw_prompt = state["user_prompt"]

    # Build a strict instruction block to prevent introspection
    instruction_block = (
        "IMPORTANT: You are only allowed to use the schema below ‚Äî you must NOT inspect or read any rows from the database. "
        "Do NOT request sample rows. Do NOT attempt to access the database for schema discovery. "
        "Using only the schema below, produce a single valid SQL query (ANSI SQL or dialect I specify if needed) that returns "
        "Return only the SQL; do not include explanation text."
    )

    combined_prompt = f"{schema_desc}\n\n{instruction_block}\n\nUser intent: {raw_prompt}\n\n"

    sql_query = vn_model.generate_sql(combined_prompt)
    #prompt = state["vanna_prompt"]

    #sql_query = vn_model.generate_sql(combined_prompt)

    try:
        result = vn_model.run_sql(sql_query)
        if isinstance(result, pd.DataFrame):
            parsed_result = result
        elif isinstance(result, list):
            parsed_result = pd.DataFrame(result)
        else:
            parsed_result = pd.DataFrame([{"Result": str(result)}])
    except Exception as e:
        parsed_result = pd.DataFrame([{"Error": f"SQL Execution failed: {e}"}])

    return {
        **prune_state(state, STATE_KEYS_SET_AT_ENTRY),
        "sql_result": parsed_result,
        "sql_query": sql_query
    }


#-----------------------------WEB SEARCH AGENT-----------------------------------
# --Enhance Google Search--
DOMAINS = {
    "core": [
        "swissre.com", "munichre.com", "amwins.com", "willistowerswatson.com",
        "insurancebusinessmag.com", "businessinsurance.com",
        "insuranceinsider.com", "iii.org",  # Insurance Information Institute
        "deloitte.com", "mckinsey.com", "bcg.com"
    ],
    "regulators": [
        "irdai.gov.in", "naic.org", "eba.europa.eu", "eiopa.europa.eu"
    ],
    "market_news": [
        "dowjones.com", "wsj.com", "reuters.com", "bloomberg.com",
        "financialtimes.com", "economist.com"
    ],
}

INSURANCE_SYNONYMS = {
    "loss ratio": ["combined ratio", "claims ratio"],
    "reserve": ["ibnr", "loss reserves", "case reserve", "ultimate loss"],
    "premium": ["written premium", "earned premium", "gross premium", "gwp"],
    "social inflation": ["nuclear verdicts", "litigation costs", "jury awards"],
}


def _domain_filter(for_news: bool) -> str:
    # Bias to relevant sources; include market_news when for_news = True
    domains = DOMAINS["core"] + DOMAINS["regulators"] + (DOMAINS["market_news"] if for_news else [])
    return " OR ".join([f"site:{d}" for d in domains])


def enhance_query(prompt: str) -> dict:
    """
    Builds query and mode.
    Returns: {"q": <string>, "for_news": bool}
    """
    p = prompt.strip()
    lower = p.lower()
#   for_news = any(w in lower for w in ["news", "today", "latest", "update", "trend", "q3", "q4", "fy", "quarter", "yoy", "benchmark"])
    for_news = "TRUE"

    insurance_tokens = ["insurance", "insurer", "claim", "premium", "underwriting", "actuarial", "reinsurance", "coverage", "reserving"]
    base_query = p if any(t in lower for t in insurance_tokens) else f"in insurance industry: {p}"
    sites = _domain_filter(for_news)

    q = f'{base_query} ({sites})'
    return {"q": q, "for_news": for_news}

# --- SerpAPI Node --- 
def serp_node(state: GraphState) -> GraphState:
    built = enhance_query(state["user_prompt"])
    query, for_news = built["q"], built["for_news"]

    search = GoogleSearch({
    "q": query,
    "api_key": os.getenv("SERPAPI_API_KEY"),
    "num": 5
    })
    results = search.get_dict()

    links = []
    summaries = []

    if "organic_results" in results:
        for r in results["organic_results"][:5]:
            link = r.get("link")
            title = r.get("title", "Untitled").strip('"')
            snippet = r.get("snippet", "No summary available.").strip('"')
            if link:
                links.append(f"[{title}]({link})")
                summaries.append(snippet)

    if not links:
        links = ["No high-quality results found (try broader query or remove filters)."]
        summaries = [""]

    # Build LLM prompt
    combined_text = "\n".join([f"[{i+1}] {s}" for i, s in enumerate(summaries)])

    sql_in_context = isinstance(state.get("sql_result"), pd.DataFrame) and not state["sql_result"].empty
    internal_sql_top5 = state["sql_result"].head(5).to_markdown(index=False) if sql_in_context else ""

    # 2) EXTERNAL: fetch market benchmark only from the web
    #    We override the user_prompt for the search node so it explicitly asks for market/industry averages externally.
    ext_prompt = f"""
    Market / industry average IBNR trend (P&C) for recent 1‚Äì5 years, external sources only.
    Prefer credible sources (Dow Jones/WSJ/Reuters/BusinessInsurance/InsuranceInsider/regulators/reinsurers).
    Use USD, %, ratios if available.
    Original user context: {state.get('user_prompt','')}
    """.strip()


    if sql_in_context:
        general_summary_prompt = f"""
        You are an insurance and actuarial analyst comparing internal company data with external web results.

        Use the following INTERNAL SQL DATA ONLY FOR CONTEXT. **Do not include internal tables or numbers in your output.**

        üßæ Internal SQL Query:
        {state['sql_query'] if 'sql_query' in state else ''}

        üìä Top 5 rows of SQL Output (reference only, do not display):
        SQL: {state.get('sql_query','')}
        Top rows:
        {internal_sql_top5}

        External snippets (numbered):
        {combined_text}
        
        User Prompt:
        "{ext_prompt}"

        üîΩ Your Task:
        - Summarize **only what is found in the external data**
        - DO NOT display the internal SQL data or repeat it
        - Be concise, no more than **6-8 lines**
        - Include **percentages, currency, loss ratios, IBNR**, and other KPIs found in the web
        - Avoid repeating full articles or sentences
        - Mention key **KPIs** (e.g., IBNR, premiums, loss ratios, reserves)
        -Focus more on numerical insights

        Output format:
        1. üìå Start with a summary of overall findings with around 5-6 lines.
        2. üî¢ Then list 6‚Äì7 **quantitative highlights**.
        3. üí¨ End with any notable quote or number from a source if applicable.
        4. Can include a table with numerical insights as well, but not the internal data or tabular data. Only if you found it in external data.
        """
    else:
        general_summary_prompt = f"""
        Your task is to extract **concise and numerically rich insights** from the following web snippets, in response to this user query:

        "{state['user_prompt']}"

        External snippets (numbered):
        {combined_text}

       Your summary should:
        - Be structured and no more than **10‚Äì12 lines**
        - Include **percentages**, **currency values**, **ratios**, **dates**, and **growth trends**
        - Mention key **KPIs** (e.g., IBNR, premiums, loss ratios, reserves)
        - Avoid repeating the snippets. Instead, **synthesize them**
        - If no numbers are found, say so explicitly

        Output format:
        1. üìå Start with a summary of overall findings with around 5-6 lines.
        2. üî¢ Then list 3‚Äì4 **quantitative highlights**.
        3. üí¨ End with any notable quote or number from a source if applicable.
        4. Can include a table with numerical insights as well
        """

    general_summary = call_llm(general_summary_prompt).strip()

    return {
        **prune_state(state, STATE_KEYS_SET_AT_ENTRY),
        "web_links": list(zip(links, summaries)),
        "general_summary": general_summary
    }


# ------------------ COMP Node  -----------------
# COMPARISON_TERMS = [
#     "market average", "industry average", "industry", "benchmark", "benchmarks",
#     "peer", "peers", "peer set", "vs", "versus", "against", "relative to",
#     "compare", "comparison", "compared to", "market trend", "external"
# ]
def _read_table(path_or_uploaded):
    """
    Robust reader for either:
      - filesystem path (str) -> use pd.read_csv / pd.read_excel
      - UploadedFile-like object (has .name and .read()) -> read bytes/text and parse
    Returns a pandas.DataFrame or raises a helpful error.
    """
    # If we were passed a path string, read directly (preferred)
    if isinstance(path_or_uploaded, str):
        path = path_or_uploaded
        if not os.path.exists(path):
            raise ValueError(f"File path does not exist: {path}")
        ext = os.path.splitext(path)[1].lower().lstrip(".")
        if ext == "csv":
            # read CSV normally (allow inference), we'll normalize keys later
            return pd.read_csv(path)
        else:
            # excel (xls/xlsx) or other spreadsheet formats
            return pd.read_excel(path)

    # If UploadedFile-like object or bytes-like
    if hasattr(path_or_uploaded, "read"):
        try:
            # Reset pointer if possible (Streamlit UploadedFile might have been read earlier)
            if hasattr(path_or_uploaded, "seek"):
                try:
                    path_or_uploaded.seek(0)
                except Exception:
                    pass

            raw = path_or_uploaded.read()
            if not raw:
                # try resetting and reading again
                if hasattr(path_or_uploaded, "seek"):
                    try:
                        path_or_uploaded.seek(0)
                        raw = path_or_uploaded.read()
                    except Exception:
                        pass

            # raw may be bytes or str
            if isinstance(raw, bytes):
                name = getattr(path_or_uploaded, "name", "") or ""
                ext = name.lower().split(".")[-1] if "." in name else None
                if ext == "csv":
                    text = raw.decode("utf-8", errors="replace")
                    return pd.read_csv(io.StringIO(text))
                else:
                    # try excel from bytes
                    try:
                        return pd.read_excel(io.BytesIO(raw))
                    except ValueError:
                        # fallback to openpyxl engine (xlsx)
                        return pd.read_excel(io.BytesIO(raw), engine="openpyxl")
            else:
                # raw is likely str (text), treat as CSV text
                text = str(raw)
                return pd.read_csv(io.StringIO(text))

        except Exception as e:
            raise ValueError(f"Failed to read uploaded file-like object: {e}")

    # Otherwise unsupported type
    raise ValueError(f"Unsupported file object: {type(path_or_uploaded)}")


def is_year_like_series(s, min_year=1900, max_year=2100, max_unique_for_key=200):
    """Return True if a pandas Series looks like a year/key column."""
    vals = s.dropna().astype(str).str.strip()
    if vals.empty:
        return False

    # If most values look like 4-digit integers within a year range -> year-like
    digit_mask = vals.str.fullmatch(r"\d{4}")
    if digit_mask.sum() / len(vals) > 0.8:
        try:
            ints = vals[digit_mask].astype(int)
            if ints.between(min_year, max_year).all():
                return True
        except Exception:
            pass

    # Also treat small-cardinality numeric columns with values in year-range as year-like
    try:
        numeric_vals = pd.to_numeric(vals, errors="coerce").dropna()
        if not numeric_vals.empty:
            unique_count = numeric_vals.nunique()
            if unique_count <= max_unique_for_key and numeric_vals.between(min_year, max_year).all():
                return True
    except Exception:
        pass

    return False


def detect_semantic_key_columns(df, extra_keywords=None):
    """
    Detect columns that should be treated as non-numeric keys even if numeric-looking.
    Returns set of column names to force as non-numeric.
    """
    if extra_keywords is None:
        extra_keywords = []

    keywords = ["year", "exposure", "exp", "underwriting", "uw", "policyyear", "policy_year", "period"]
    keywords = keywords + extra_keywords
    keywords = [k.lower() for k in keywords]

    forced_keys = set()
    cols = list(df.columns)

    # 1) name-based detection: regex 'year' or close match tokens
    for col in cols:
        col_low = str(col).lower()
        if any(re.search(rf"\b{k}\b", col_low) for k in keywords):
            forced_keys.add(col)
            continue

        tokens = re.split(r"[^a-z0-9]+", col_low)
        for t in tokens:
            if t and get_close_matches(t, keywords, n=1, cutoff=0.8):
                forced_keys.add(col)
                break

    # 2) value-based detection (fallback): numeric columns that look like years
    for col in cols:
        if col in forced_keys:
            continue
        try:
            if is_year_like_series(df[col]):
                forced_keys.add(col)
        except Exception:
            continue

    return forced_keys


def normalize_key_column(df, col):
    """Normalize a column intended for use as a join key (in-place)."""
    if col not in df.columns:
        return

    s = df[col]
    # If numeric dtype (int/float), convert integer-like floats to ints then to str
    if pd.api.types.is_numeric_dtype(s):
        nonnull = s.dropna()
        try:
            if not nonnull.empty and np.allclose(nonnull, np.round(nonnull)):
                # integer-like floats or ints -> safe to cast to int then str
                df[col] = s.where(s.isna(), s.round(0).astype(int).astype(str))
            else:
                # keep numeric but convert to string removing trailing .0
                df[col] = s.astype(str).str.replace(r'\.0+$', '', regex=True).str.strip()
        except Exception:
            df[col] = s.astype(str).str.replace(r'\.0+$', '', regex=True).str.strip()
    else:
        # object/string dtype: remove trailing .0, commas, extra whitespace
        s2 = s.astype(str).str.strip()
        s2 = s2.str.replace(r'\.0+$', '', regex=True)
        s2 = s2.str.replace(',', '', regex=False).str.strip()
        s2 = s2.str.replace(r'\s+', ' ', regex=True)
        df[col] = s2.replace({"nan": np.nan, "None": np.nan})


def normalize_key_columns_for_join(df, extra_keywords=None):
    """Auto-detect and normalize key-like columns in df (in-place)."""
    # Generic trim for all object columns
    for c in df.columns:
        if pd.api.types.is_object_dtype(df[c]):
            df[c] = df[c].astype(str).str.strip()

    # Detect candidate semantic key columns
    forced = detect_semantic_key_columns(df, extra_keywords=extra_keywords)
    for col in forced:
        normalize_key_column(df, col)

    # Additional standardization for common categorical columns (optional)
    # Title-case UW Unit-style columns and strip RI Type strings
    for ucol in df.columns:
        if ucol and isinstance(ucol, str) and re.search(r"uw|underwrit", ucol, re.IGNORECASE):
            df[ucol] = df[ucol].astype(str).str.strip().str.replace(r'\s+', ' ', regex=True).str.title()
    for rcol in df.columns:
        if rcol and isinstance(rcol, str) and re.search(r"ri|reinsur|re-insur|ri type", rcol, re.IGNORECASE):
            df[rcol] = df[rcol].astype(str).str.strip().str.replace(r'\s+', ' ', regex=True)

    return forced  # return the set of normalized columns


# --- Apply accounting / presentation formatting ---
def accounting_format(x):
    """Format numeric values in accounting style: commas, 2 decimals, parentheses for negatives."""
    try:
        if pd.isna(x):
            return ""
        if isinstance(x, (int, float, np.integer, np.floating)):
            return f"({abs(x):,.2f})" if x < 0 else f"{x:,.2f}"
        return x
    except Exception:
        return x

def reconciliation_node(file1, file2):
    # 1) Read inputs
    df_A = _read_table(file1)
    df_B = _read_table(file2)

    # 2) Normalize header names (strip whitespace)
    df_A.columns = [c.strip() if isinstance(c, str) else c for c in df_A.columns]
    df_B.columns = [c.strip() if isinstance(c, str) else c for c in df_B.columns]

    # 3) Detect semantic key columns and normalize them in-place
    forced_A = normalize_key_columns_for_join(df_A, extra_keywords=["expo", "exposureyear", "exp_year"])
    forced_B = normalize_key_columns_for_join(df_B, extra_keywords=["expo", "exposureyear", "exp_year"])

    # Union of forced columns from both sides ‚Äî ensure both sides normalized for the same keys
    forced_union = set(forced_A) | set(forced_B)
    for col in forced_union:
        if col in df_A.columns:
            normalize_key_column(df_A, col)
        if col in df_B.columns:
            normalize_key_column(df_B, col)

    # 4) Determine numeric vs non-numeric columns (force union keys into non-numeric)
    numeric_cols = []
    non_numeric_cols = []
    for col in df_A.columns:
        if col in forced_union:
            non_numeric_cols.append(col)
            continue
        # treat a column as numeric only if pandas considers it numeric
        if pd.api.types.is_numeric_dtype(df_A[col]):
            numeric_cols.append(col)
        else:
            non_numeric_cols.append(col)

    # Safety: ensure we have at least one non-numeric key to join on
    if len(non_numeric_cols) == 0:
        raise ValueError("No non-numeric key columns detected to join on. Please ensure there are categorical key columns (e.g., UW Unit, Exposure Year, RI Type).")

    # 5) Before merging, canonicalize non-numeric columns to string on both dfs to avoid dtype mismatches
    # === normalize falsey placeholder strings into real NaN ===
    pattern = r'^\s*(none|null|nan|na|nan\.0+)?\s*$'  # case-insensitive match for common empty-like tokens
    for col in non_numeric_cols:
        if col in df_A.columns:
            df_A[col] = df_A[col].astype(str).str.strip()
            # mask rows where the cell matches the placeholder pattern (case-insensitive)
            mask_a = df_A[col].str.match(pattern, case=False, na=False)
            df_A.loc[mask_a, col] = np.nan
            # also convert empty strings to np.nan
            df_A[col].replace('', np.nan, inplace=True)

        if col in df_B.columns:
            df_B[col] = df_B[col].astype(str).str.strip()
            mask_b = df_B[col].str.match(pattern, case=False, na=False)
            df_B.loc[mask_b, col] = np.nan
            df_B[col].replace('', np.nan, inplace=True)

    # Now drop source rows where ALL key columns are missing (invalid combos)
    df_A = df_A[~df_A[non_numeric_cols].isna().all(axis=1)].copy()
    df_B = df_B[~df_B[non_numeric_cols].isna().all(axis=1)].copy()

    # === outer merge ===
    merged = pd.merge(df_A, df_B, on=non_numeric_cols, how="outer", suffixes=("_A", "_B"))

    # Defensive: remove merged rows where ALL join keys are missing (just in case)
    if not merged.empty:
        all_key_null_mask = merged[non_numeric_cols].isna().all(axis=1)
        if all_key_null_mask.any():
            # optional debug logging (uncomment in dev): print or store sample rows
            # st.write("Dropping merged rows with empty keys:", merged.loc[all_key_null_mask, non_numeric_cols].head().to_dict(orient="records"))
            merged = merged.loc[~all_key_null_mask].reset_index(drop=True)

    # === compute missing lists ===
    a_mask = merged.filter(like="_A").isna().any(axis=1)
    b_mask = merged.filter(like="_B").isna().any(axis=1)

    missing_in_A = merged.loc[a_mask, non_numeric_cols].drop_duplicates().reset_index(drop=True)
    missing_in_B = merged.loc[b_mask, non_numeric_cols].drop_duplicates().reset_index(drop=True)

    # Final presentation cleanup: drop any rows in missing lists that still contain NaN in any key
    missing_in_A = missing_in_A.dropna(subset=non_numeric_cols).reset_index(drop=True)
    missing_in_B = missing_in_B.dropna(subset=non_numeric_cols).reset_index(drop=True)


    # Variance computation
    variance_df = []
    for col in numeric_cols:
        merged[f"{col}_abs_diff"] = merged[f"{col}_A"] - merged[f"{col}_B"]
        merged[f"{col}_pct_diff"] = ((merged[f"{col}_abs_diff"]) / (merged[f"{col}_B"].replace(0, np.nan))) * 100
    # Build the variance DataFrame (includes non-numeric keys + all diff columns)
    variance_df = merged[non_numeric_cols + [c for c in merged.columns if "_diff" in c]].copy()

    # --- Keep only rows where at least one diff column is non-zero ---
    diff_cols = [c for c in variance_df.columns if "_diff" in c]

    # Drop rows where all diff columns are either NaN or 0
    if diff_cols:
        variance_df = variance_df.loc[
            ~(variance_df[diff_cols].fillna(0).apply(lambda x: np.allclose(x, 0), axis=1))
        ].reset_index(drop=True)

    # Format all numeric diff columns
    for col in diff_cols:
        variance_df[col] = variance_df[col].apply(accounting_format)


    # Variance summary
    summary = {
        "total_variance": merged[[f"{col}_abs_diff" for col in numeric_cols]].sum().to_dict(),
        "avg_pct_diff": merged[[f"{col}_pct_diff" for col in numeric_cols]].mean().to_dict(),
    }


    # --- Prepare compact missing summaries for LLM ---
    def summarize_missing(df, label):
        """Convert a missing-combination DataFrame into a short human-readable list."""
        if df.empty:
            return f"No missing combinations in {label}."
        # Convert to strings of key tuples for readability
        sample_records = df.head(10).to_dict(orient="records")
        text = ", ".join(
            [" | ".join(str(v) for v in rec.values()) for rec in sample_records]
        )
        more_note = "" if len(df) <= 10 else f" ... and {len(df)-10} more."
        return f"{len(df)} missing combinations in {label}: {text}{more_note}"

    missing_summary_A = summarize_missing(missing_in_A, "Dataset A")
    missing_summary_B = summarize_missing(missing_in_B, "Dataset B")

    # Variance commentary (LLM)
    prompt = f"""
    You are a financial analyst reviewing quarterly reconciliation results.

    Compare two Excel datasets and write a concise, professional financial commentary summarizinge:
    - Overall variance across numeric fields
    - Key increases/decreases
    - Missing combinations in A and B

    Summary stats:
    {summary}

    Sample variances:
    {variance_df.head(5).to_dict(orient='records')}

    Missing combination overview:
    Missing combinations in A: {missing_summary_A}
    Missing combinations in B: {missing_summary_B}
    """
    variance_commentary = call_llm(prompt)  # call your LLM here

    return {
        "variance_df": variance_df,
        "missing_in_A": missing_in_A,
        "missing_in_B": missing_in_B,
        "variance_summary": summary,
        "variance_commentary": variance_commentary,
    }




def comp_node(state: GraphState) -> GraphState:

    file1 = state.get("uploaded_file1_path")
    file2 = state.get("uploaded_file2_path")

    # --- Determine Comp Node Type ---
    if file1 and file2:
        comp_type = "EXCEL-EXCEL"
    else:
        comp_type = "SQL-WEB"

    # If 2 excels are attached ‚Üí reconciliation flow
    if file1 and file2:
        reconciliation_result = reconciliation_node(file1, file2)
        state["reconcile_df"] = reconciliation_result["variance_df"]
        state["missing_in_A"] = reconciliation_result["missing_in_A"]
        state["missing_in_B"] = reconciliation_result["missing_in_B"]
        state["variance_summary"] = reconciliation_result["variance_summary"]
        state["variance_commentary"] = reconciliation_result["variance_commentary"]
        state["comp_type"] = comp_type
        return state

    else:
        # 1) INTERNAL: build a safe Vanna prompt (internal-only)

        schema_desc = get_schema_description('Actuarial_Data.db')
        raw_prompt = state.get("vanna_prompt") or state["user_prompt"]

        # Build a strict instruction block to prevent introspection
        instruction_block = (
            "IMPORTANT: You are only allowed to use the schema below ‚Äî you must NOT inspect or read any rows from the database. "
            "Do NOT request sample rows. Do NOT attempt to access the database for schema discovery. "
            "Using only the schema below, produce a single valid SQL query (ANSI SQL or dialect I specify if needed) that returns "
            "Return only the SQL; do not include explanation text."
        )

        combined_prompt = f"{schema_desc}\n\n{instruction_block}\n\nUser intent: {raw_prompt}\n\n"

        sql_query = vn_model.generate_sql(combined_prompt)

        try:
            result = vn_model.run_sql(sql_query)
            if isinstance(result, pd.DataFrame):
                sql_df = result
            elif isinstance(result, list):
                sql_df = pd.DataFrame(result)
            else:
                sql_df = pd.DataFrame([{"Result": str(result)}])
        except Exception as e:
            sql_df = pd.DataFrame([{"Error": f"SQL Execution failed: {e}"}])

        serp_result = serp_node({**state, "sql_query": sql_query, "sql_result": sql_df})

        web_links = serp_result.get("web_links", [])
        external_summary = serp_result.get("general_summary", "")

        # 3) COMPARISON: clear separation + citations
        comparison_prompt = f"""
        You are an Benchmarking actuarial analyst. Compare OUR internal IBNR trend to EXTERNAL market/industry benchmarks.
        Rules:
        - Use INTERNAL SQL only for our numbers; do NOT infer market values from internal data.
        - Use EXTERNAL WEB snippets only for market/industry values; if no numeric market average is found, say so explicitly.
        - Put all money in **USD**, include **%/ratios/dates** where present.
        - Append [i] citations for any external metric where i refers to the snippet index (shown below).
        - If sources disagree, note the discrepancy briefly.

        Your job is to:
        1. Analyze differences, similarities, and gaps between internal company data and external web sources.
        2. Focus heavily on **numerical metrics** such as:
        - IBNR, Incurred Loss, Ultimate Loss
        - Premiums, Loss Ratios
        - Exposure Years, Percent changes

        3. Focus more on:
        - Trends (increases/decreases)
        - Matching vs. diverging figures
        - Numerical differences or % differences

        Our internal (context only; do not reveal raw table):
        SQL: {sql_query}
        Top rows (context only):
        {sql_df.head(5).to_markdown(index=False) if isinstance(sql_df, pd.DataFrame) else str(sql_df)}

        External snippets (numbered):
        {chr(10).join([f"[{i+1}] {s}" for i, (_, s) in enumerate(web_links)])}

        Task:
        1) 5‚Äì7 lines overview (internal vs market).
        2) 3‚Äì5 bullets with side-by-side contrasts (Our vs Market) using USD/%/ratios and [citation] only for external numbers.
        3) 1 ‚Äúwatch item‚Äù (e.g., social inflation, rate adequacy, reserving pressure) if relevant.

        General external synthesis to leverage (do not copy verbatim; keep citations): 
        {external_summary}
        """
        comparison_summary = call_llm(comparison_prompt).strip()

        return {
            **prune_state(state, STATE_KEYS_SET_AT_ENTRY),
            "sql_result": sql_df,
            "sql_query": sql_query,
            "web_links": web_links,
            "general_summary": external_summary,
            "comparison_summary": comparison_summary,
            "comp_type": comp_type
        }



# faissdb node to extract internal docs
def faissdb_node(state: GraphState) -> GraphState:
    faiss = FAISS.load_local(
        folder_path="faiss_index/",
        embeddings=OpenAIEmbeddings(),
        allow_dangerous_deserialization=True
    )
    docs = faiss_index.similarity_search(state["user_prompt"], k=3)

    top_docs = docs[:3]  # ‚¨ÖÔ∏è Top 3 instead of 5
    content_snippets = "\n\n---\n\n".join(d.page_content[:500] for d in top_docs)

    summary_prompt = f"""
    Based on the following retrieved document chunks from internal knowledge base, answer the user's query:

    User Prompt: {state['user_prompt']}

    Documents:
    {content_snippets}

    Provide a concise and structured answer with key points or numeric details if found.
    """
    summary = call_llm(summary_prompt)

    # Extract faiss_sources with source path
    faiss_sources = []
    all_images = []

    for doc in top_docs:
        doc_name = doc.metadata.get("source_doc", "Doc")
        snippet = doc.page_content[:300]
        path = doc.metadata.get("file_path")  # must be present in ingestion step
        #print(f"[DEBUG] FAISS doc metadata: {doc.metadata}")
        faiss_sources.append((doc_name, snippet, path))

        # Load associated images
        image_meta_path = os.path.join("extracted_images", "extracted_image_metadata.json")
        if os.path.exists(image_meta_path):
            with open(image_meta_path, 'r') as f:
                all_metadata = json.load(f)
            related_images = [
                meta for meta in all_metadata
                if meta["original_doc"] == doc_name
            ]
            all_images.extend(related_images)

    return {
        **prune_state(state, STATE_KEYS_SET_AT_ENTRY),
        "faiss_summary": summary,
        "faiss_sources": faiss_sources,
        "faiss_images": all_images
    }



#-------------------------DOCUMENT NODE-------------------------------

# ----------------------------
# ü§ñ Ollama Fuzzy Header Match
# ----------------------------
def get_target_header_and_table(state: GraphState) -> GraphState:

    """
    Uses OpenAI LLM to identify the most relevant header and table from the document,
    based on the user's instruction (fuzzy_prompt or user_prompt).
    Updates GraphState with header_candidate and table_candidate_index.
    """
    

    doc_path = state.get("document_path")
    if not doc_path:
        st.error("‚ùå No document path found in state.")
        return state

    # Load doc and extract structure
    doc = Document(doc_path)
    structure = extract_structure(doc)
    structure_str = stringify_structure(structure)

    # Choose prompt (prefer fuzzy_prompt if available)
    instruction = state.get("user_prompt")

    prompt = f"""
    You are helping identify the correct table to update in a Word document.

    The document contains several sections. Each section has a header and a list of tables.
    Each table has:
    - a table index (starting from 0 under that header),
    - the number of rows and columns,
    - and the list of column headers (which might vary slightly across documents).

    Here is the document structure:
    {structure_str}

    The user's instruction is:
    \"\"\"{instruction}\"\"\"

    Rules:
    - Match the most relevant headers and tables to the instruction (fuzzy allowed).
    - User's instruction might have columns mentioned as well along with headers. 
      SO use them to find the exact table in case there are multiple tables under the same header.
    - If multiple tables exist under a header, rank them by column schema similarity. And fetch the most similar one.

    Return a single JSON object exactly in this form:
    {{"header_text": "Exact header from document", "table_index_under_header": 0}}

    If no good match exists, return {{}}.
    """

    try:
        response = call_llm(prompt)
        st.write(response)
        match = re.search(r"\{.*\}", response, re.DOTALL)
        if not match:
            # no parseable JSON
            return state
        parsed = json.loads(match.group())
        if not parsed:
            return state

        header_text = parsed.get("header_text")
        table_idx = parsed.get("table_index_under_header")

        # validate
        if header_text is None or table_idx is None:
            return state

        # set in state
        state["header_candidate"] = header_text
        state["table_candidate_index"] = int(table_idx)
        #state["candidate_tables"] = [{"header_text": header_text, "table_index_under_header": int(table_idx)}]
        return state

    except Exception:
        # on any failure, return state unchanged
        return state

# ----------------------------
# üìÑ Document Parsing
# ----------------------------
def extract_structure(doc: "Document") -> list:
    """
    Walk doc.element.body and return a list of sections:
      [{"header": <str>, "tables": [Table, ...]}, ...]
    If a table appears before the first heading, its section header will be "NO_HEADER".
    """
    structure = []
    current_header = None

    for element in doc.element.body:
        if isinstance(element, CT_P):
            para = Paragraph(element, doc)
            # guard for paragraphs without style
            style_name = None
            try:
                style_name = para.style.name if para.style else None
            except Exception:
                style_name = None
            if style_name and style_name.startswith("Heading"):
                current_header = para.text.strip()
                structure.append({"header": current_header, "tables": []})
        elif isinstance(element, CT_Tbl):
            table = Table(element, doc)
            if current_header is None:
                current_header = "NO_HEADER"
                structure.append({"header": current_header, "tables": []})
            structure[-1]["tables"].append(table)

    return structure

def stringify_structure(structure: list, max_chars: int = 3000) -> str:
    """
    Turn `structure` into a compact text describing each header and its tables:
    - header text
    - number of tables
    - for each table: rows, cols, first-row headers (up to 6)
    """
    summary = []
    for sec in structure:
        header = sec.get("header", "")
        tables = sec.get("tables", [])
        sec_lines = [f"HEADER: {header} (tables: {len(tables)})"]
        for ti, tbl in enumerate(tables):
            try:
                rows = len(tbl.rows)
                cols = len(tbl.columns)
                # grab first row texts as candidate column names (if exists)
                first_row = []
                if rows > 0:
                    first_row = [cell.text.strip().replace("\n", " ")[:60] for cell in tbl.rows[0].cells]
                sec_lines.append(f"  - Table {ti}: rows={rows}, cols={cols}, cols_sample={first_row[:6]}")
            except Exception:
                sec_lines.append(f"  - Table {ti}: (could not introspect)")
        summary.append("\n".join(sec_lines))

    out = "\n\n".join(summary)
    if len(out) > max_chars:
        return out[:max_chars] + " ... (truncated)"
    return out


# ----------------------------
# üîÅ Replace Table with Formatting
# ----------------------------
def get_column_widths(table: "Table") -> list:
    """
    Return a list of column widths (raw values) if available, else None entries.
    Works by inspecting table._element XML. This is best-effort.
    """
    widths = []
    try:
        # each tblGrid -> gridCol elements have 'w:w' attributes in twips
        tbl = table._element
        gridCols = tbl.xpath(".//w:tblGrid//w:gridCol", namespaces=tbl.nsmap)
        if gridCols:
            for gc in gridCols:
                w = gc.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}w")
                if w:
                    widths.append(int(w))
                else:
                    widths.append(None)
            return widths
    except Exception:
        pass

    # fallback: try to estimate equal widths
    try:
        ncols = len(table.columns)
        if ncols > 0:
            return [None] * ncols
    except Exception:
        return []

    return []


def replace_table(old_table: DocxTable, new_df: pd.DataFrame) -> DocxTable:
    """
    Replace `old_table` with a new table built from `new_df`, preserving style and column widths.

    Approach:
    1. Compute parent element and index of old table (xml).
    2. Create a new table using the same python-docx container (old_table._parent.add_table).
       That new table will be appended to the end of the container.
    3. Insert the new table's XML element into the parent's children right AFTER the old table element.
    4. Remove the old table's xml element.
    5. Return the new python-docx Table object.

    This avoids placing the new table before the old one (which is what happens if you insert at the same index then
    remove the old element immediately).
    """
    # sanity
    if not isinstance(new_df, pd.DataFrame):
        new_df = pd.DataFrame(new_df)

    parent = old_table._element.getparent()        # lxml element containing the table element
    index = parent.index(old_table._element)      # index of old_table in parent children
    widths = get_column_widths(old_table)         # best-effort widths
    style = None
    try:
        style = old_table.style
    except Exception:
        style = None

    # Create new table under same container (old_table._parent is a python-docx object that supports add_table)
    try:
        container = old_table._parent
        new_table = container.add_table(rows=1, cols=len(new_df.columns))
    except Exception as e:
        raise RuntimeError(f"Could not create new table in same container: {e}")

    # Apply style if available
    if style:
        try:
            new_table.style = style
        except Exception:
            pass

    # Fill header row
    for col_idx, col_name in enumerate(new_df.columns):
        try:
            new_table.cell(0, col_idx).text = str(col_name)
        except Exception:
            # ensure there are enough cells
            pass

    # Fill data rows
    for _, row in new_df.iterrows():
        cells = new_table.add_row().cells
        for i, val in enumerate(row):
            try:
                cells[i].text = "" if pd.isnull(val) else str(val)
            except Exception:
                pass

    # Apply column widths (best-effort)
    for i, col in enumerate(new_table.columns):
        if i < len(widths) and widths[i]:
            for cell in col.cells:
                try:
                    cell.width = widths[i]
                except Exception:
                    pass

    # Insert the new_table element into the parent right AFTER the old table element
    # (so new_table is placed immediately below the old table, then we remove the old table)
    parent.insert(index + 1, new_table._element)

    # Now remove the old table element
    parent.remove(old_table._element)

    return new_table


def render_doc_results(state: GraphState):
    st.write("DEBUG: ENTER render_doc_results")
    st.write("DEBUG: state keys in render:", list(state.keys()))
    st.subheader("üìÑ Document Update Summary (render_doc_results)")

    header = state.get("header_updated")
    table_idx = state.get("table_index_updated")

    st.markdown(f"**Header:** {header}")
    st.markdown(f"**Table index:** {table_idx}")

    if state.get("sql_query"):
        st.code(state["sql_query"], language="sql")

    sql_df = state.get("sql_result")
    st.write("DEBUG: sql_df type in render:", type(sql_df))
    if isinstance(sql_df, pd.DataFrame):
        st.dataframe(sql_df)
    else:
        st.write(sql_df)

    if state.get("updated_doc_path"):
        st.write("DEBUG: updated_doc_path exists:", state["updated_doc_path"])
        try:
            with open(state["updated_doc_path"], "rb") as f:
                st.download_button("üì• Download Updated Document", f, file_name="updated.docx")
        except Exception as e:
            st.write("DEBUG: failed to open updated_doc_path:", e)



def run_sql_and_update_doc(
    state: GraphState,
    user_prompt: str,
    old_table: Table,
    header: str,
    table_idx: int,
    doc: Document
) -> GraphState:
    """
    Generates a Vanna prompt, runs SQL, replaces the table in the document,
    saves updated doc, and returns updated GraphState.
    """

    # --- Step 1: Build Vanna Prompt using LLM ---
    col_headers = [cell.text.strip() for cell in old_table.rows[0].cells]
    st.write("col_headers")
    st.write(col_headers)
    schema = get_schema_description('Actuarial_Data.db')

    #The user has asked: "{user_prompt}". 
    #User's instruction might have columns mentioned as well along with headers. Ignore those column names.

    vanna_prompt_instruction = f"""
    You are an assistant that generates a natural language prompt for Vanna AI that Vanna will convert to SQL query in future.


    The table extracted from the Word document has these columns:
    {', '.join(col_headers)}.

    Please generate a clear SQL data retrieval prompt in natural language that include all the columns found in the word document explicitly.
    Refer to the below mentioned stuff for sql database schemas, columns and for better accuracy
    - database schema:
    {schema}
    - Use this additional documentation to better understand column meanings:
    {documentation}
    - Additionally, here are some examples of SQL-style questions and their corresponding queries (QSPairs):
    {qs_examples}

    -Example: if the headers found from the word document table are "Exposure Year", "Incurred Loss", "Ultimate Loss", 
    then your output shoud be "Show exposure year wise incurred loss and ultimate loss"
    - And always try to group by on possible dimension columns or non-numeric columns.
    """

    #st.write(vanna_prompt_instruction)
    try:
        vanna_prompt = call_llm(vanna_prompt_instruction)
        st.write("vanna_prompt")
        st.write(vanna_prompt)
    except Exception as e:
        raise RuntimeError(f"call_llm failed: {e}")

    state["vanna_prompt"] = vanna_prompt

    # --- Step 2: Generate SQL + Run ---
    try:
        sql_query = vn_model.generate_sql(vanna_prompt)
        st.write(sql_query)
        result = vn_model.run_sql(sql_query)
        new_df = pd.DataFrame(result) if not isinstance(result, pd.DataFrame) else result
    except Exception as e:
        st.error(f"‚ùå SQL execution failed: {e}")
        return state

    # --- Step 3: Replace Table in Doc ---
    try:
        replace_table(old_table, new_df)
    except Exception as e:
        raise RuntimeError(f"replace_table failed: {e}")

    # ---Step 4: save updated doc to temp file
    tmp_path = os.path.join(tempfile.gettempdir(), f"updated_{uuid.uuid4().hex}.docx")
    try:
        doc.save(tmp_path)
        download_key = f"download_updated_{uuid.uuid4().hex}"
        state["updated_doc_key"] = download_key
    except Exception as e:
        raise RuntimeError(f"Saving updated document failed: {e}")
    
    # --- Step 4: Update State ---
    state = {
        **state,
        "sql_result": new_df,
        "sql_query": sql_query,
        "updated_doc_path": tmp_path,
        "header_updated": header,
        "table_index_updated": table_idx,
        "vanna_prompt": vanna_prompt
    }

    return state


def document_node(state: GraphState) -> GraphState:
    """
    Automated document node (single-candidate flow).
    Uses existing helpers:
      - get_target_header_and_table(state) -> should set header_candidate & table_candidate_index
      - extract_structure(doc)
      - stringify_structure(structure)
      - run_sql_and_update_doc(state, user_prompt, old_table, header, table_idx, doc)
      - replace_table(old_table, new_df) (used by run_sql_and_update_doc)
    Flow:
      1. Load persisted graph_state
      2. Call fuzzy match to populate header_candidate/table_candidate_index
      3. Locate the table object, show preview
      4. Run run_sql_and_update_doc (performs vanna SQL + replace + save)
      5. Persist state and render final outputs + download button
    """

    # 0. Prefer persisted graph_state if present (so results survive reruns)
    persisted = st.session_state.get("graph_state")
    if persisted:
        # merge persisted over incoming state so persisted values have precedence
        state = {**state, **persisted}

    # Basic checks
    user_prompt = state.get("user_prompt")
    doc_path = state.get("document_path")

    if not user_prompt:
        st.error("No user prompt found. Submit the query form first.")
        return state
    if not doc_path:
        st.error("No document uploaded. Upload and submit the form first.")
        return state

    # # If already ran earlier and results exist, just render them and return
    # if state.get("sql_result") is not None:
    #     # Render summary + results + download button
    #     #st.success("‚úÖ Document has already been updated. Showing saved results.")
    #     header = state.get("header_updated") or state.get("header_candidate")
    #     table_idx = state.get("table_index_updated") or state.get("table_candidate_index")

    #     if header:
    #         st.markdown(f"**Header updated / matched:** {header}")
    #     if table_idx is not None:
    #         st.markdown(f"**Table index:** {table_idx}")

    #     if state.get("sql_query"):
    #         st.subheader("SQL Query")
    #         st.code(state["sql_query"], language="sql")

    #     st.subheader("SQL Result")
    #     sql_df = state.get("sql_result")
    #     if isinstance(sql_df, pd.DataFrame):
    #         st.dataframe(sql_df)
    #     else:
    #         st.write(sql_df)

    #     updated_path = state.get("updated_doc_path")
    #     if updated_path and os.path.exists(updated_path):
    #         with open(updated_path, "rb") as f:
    #             st.download_button("üì• Download Updated Document", f, file_name="updated.docx")
    #     else:
    #         st.warning("Updated document not available (temp file may have been removed).")

    #     return state


    # 1. Load document & extract structure
    doc = Document(doc_path)
    structure = extract_structure(doc)
    structure_str = stringify_structure(structure)

    # 2. Get single best match (this helper should set header_candidate & table_candidate_index in state)
    state = get_target_header_and_table(state)
    header = state.get("header_candidate")
    table_idx = state.get("table_candidate_index")

    if not header or table_idx is None:
        st.error("Could not automatically identify a matching header/table in the document.")
        # optionally show the structure for debugging
        st.subheader("Document structure (debug)")
        st.text(structure_str[:2000] if structure_str else "No structure")
        return state

    # 3. Locate the Table object in the structure
    try:
        target_section = next(sec for sec in structure if sec["header"] == header)
        old_table = target_section["tables"][table_idx]
        
    except Exception as e:
        st.error(f"Failed to locate the table found by the fuzzy matcher: {e}")
        return state

    # 4. Show the matched table preview
    #st.subheader("Matched table (preview)")
    preview_df = pd.DataFrame([[cell.text for cell in row.cells] for row in old_table.rows])
    #st.dataframe(preview_df)

    # Store a serializable representation in state (list of dicts)
    state["preview_df"] = preview_df.to_dict(orient="records")

    # If you also want column names separately:
    state["preview_df_columns"] = list(preview_df.columns)

    
    # 5. Run SQL + update doc automatically (show spinner while running)
    #st.info("Generating Vanna prompt, creating SQL, executing and updating the document...")
    #with st.spinner("Running LLM/Vanna and updating document..."):
    try:
        # run_sql_and_update_doc should:
        # - generate vanna_prompt (LLM)
        # - generate & run SQL (vn_model)
        # - replace_table(old_table, new_df)
        # - save updated document and return updated state with sql_result, sql_query, updated_doc_path, header_updated, table_index_updated, vanna_prompt
        new_state = run_sql_and_update_doc(state, user_prompt, old_table, header, table_idx, doc)

    except Exception as e:
        st.error(f"Automated update failed: {e}")
        # optionally show traceback for debugging during development
        try:
            import traceback
            st.text(traceback.format_exc())
        except Exception:
            pass
        return state

    # 6. Persist the returned state and render final outputs
    # Ensure we got expected keys back
    state = {**state, **(new_state or {})}
    st.session_state.graph_state = state

    #st.success("‚úÖ Document updated successfully.")
    header = state.get("header_updated") or header
    table_idx = state.get("table_index_updated") or table_idx

    #st.markdown(f"**Header updated:** {header}")
    #st.markdown(f"**Table index updated:** {table_idx}")

    #if state.get("sql_query"):
    #    st.subheader("SQL Query")
    #    st.code(state["sql_query"], language="sql")

    #st.subheader("SQL Result")
    #sql_df = state.get("sql_result")
    #if isinstance(sql_df, pd.DataFrame):
    #    st.dataframe(sql_df)
    #else:
    #    st.write(sql_df)

#    updated_path = state.get("updated_doc_path")
#    if updated_path and os.path.exists(updated_path):
#        with open(updated_path, "rb") as f:
#            st.download_button("üì• Download Updated Document", f, file_name="updated.docx", key=state.get("updated_doc_key"))
#    else:
#        st.warning("Updated document not available (temp file may have been removed).")

    return state


def generate_follow_up_questions(user_prompt: str) -> List[str]:
    followup_prompt = f"""
    Based on the following insurance-related user query:
    "{user_prompt}"

    Suggest 3 intelligent follow-up questions the user could ask next. Keep them short, relevant, and not repetitive.
    Return them as a plain list.
    """
    try:
        response = call_llm(followup_prompt)
        return re.findall(r"^\s*[-‚Äì‚Ä¢]?\s*(.+)", response, re.MULTILINE)[:3] or response.split("\n")[:3]
    except:
        return []



def visualize_workflow(builder, active_route=None):

    route_to_node = {
        "sql": "vanna_sql",
        "search": "serp_search",
        "document": "doc_update",
        "faissdb": "faissdb",
        "comp": "comp"
    }

    highlight_node = route_to_node.get(active_route)

    G = nx.DiGraph()
    edge_styles = {}

    # Add all nodes
    for node in builder.nodes:
        G.add_node(node)
    G.add_node("__start__")
    G.add_node("__end__")

    # LangGraph base edges
    for source, target in builder.edges:
        G.add_edge(source, target)
        # Add style only for non-router edges
        if source != "router":
            edge_styles[(source, target)] = {"style": "solid", "color": "black", "width": 1.5}

    # Always show dashed edges from router to all 3 branches
    for target in ["vanna_sql", "serp_search", "doc_update", "comp", "faissdb"]:
        if ("router", target) not in G.edges:
            G.add_edge("router", target)
        edge_styles[("router", target)] = {"style": "dashed", "color": "gray", "width": 1}

    # Highlight the active route in red
    if highlight_node:
        edge_styles[("router", highlight_node)] = {"style": "solid", "color": "red", "width": 2.5}

    # Positions for nodes
    pos = {
        "__start__": (2, 4),
        "router": (2, 3),
        "vanna_sql": (0, 2),
        "serp_search": (1, 2),
        "doc_update": (2, 2),
        "comp": (3, 2),
        "faissdb": (4, 2),
        "__end__": (2, 1),
    }

    plt.figure(figsize=(6, 5))
    nx.draw_networkx_nodes(G, pos, node_size=2500, node_color="skyblue")
    nx.draw_networkx_labels(G, pos, font_size=10, font_weight="bold")

    # Draw styled edges
    for edge in G.edges:
        style = edge_styles.get(edge, {"style": "solid", "color": "black", "width": 1})
        nx.draw_networkx_edges(
            G, pos,
            edgelist=[edge],
            arrows=True,
            arrowstyle='-|>',
            style=style["style"],
            edge_color=style["color"],
            width=style["width"]
        )

    plt.title("Agentic LangGraph Workflow")
    plt.axis("off")
    plt.tight_layout()
    st.pyplot(plt)


def _rows_cols_from_serialized(df_like):
    """
    Accepts:
      - pandas.DataFrame
      - dict with {"columns": [...], "rows": [...]}
      - list[dict] (rows only)
    Returns: (columns:list[str], rows:list[list[str]])
    """
    if df_like is None:
        return [], []
    # DataFrame
    if isinstance(df_like, pd.DataFrame):
        cols = list(df_like.columns)
        rows = df_like.to_dict(orient="records")
        return cols, [[str(row.get(c, "")) for c in cols] for row in rows]
    # {"columns": [...], "rows": [...]}
    if isinstance(df_like, dict) and "rows" in df_like:
        cols = df_like.get("columns") or []
        rows_data = df_like["rows"]
        # if columns missing, infer from first row
        if not cols and isinstance(rows_data, list) and rows_data:
            cols = list(rows_data[0].keys())
        rows = []
        for r in rows_data or []:
            if isinstance(r, dict):
                rows.append([str(r.get(c, "")) for c in cols])
            else:
                # row already list-like
                rows.append([str(v) for v in (r or [])])
        return cols, rows
    # list-of-dicts
    if isinstance(df_like, list) and (not df_like or isinstance(df_like[0], dict)):
        cols = list(df_like[0].keys()) if df_like else []
        rows = [[str(r.get(c, "")) for c in cols] for r in df_like]
        return cols, rows
    # Fallback: treat as string
    return ["value"], [[str(df_like)]]


def _add_table_slide(prs, title, columns, rows, max_rows=6):
    """
    Adds a slide with a table safely.
    - columns: list[str] or []
    - rows: list[list[str]] where each row may have variable length
    - max_rows: maximum number of data rows to show (headers + data <= max_rows)
    """
    layout = prs.slide_layouts[5]  # title + content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    # Normalize rows to list of lists of strings
    norm_rows = []
    for r in rows or []:
        if isinstance(r, dict):
            # if row is dict, map by columns if available
            if columns:
                norm_rows.append([str(r.get(c, "")) for c in columns])
            else:
                # dict -> preserve order of keys
                norm_rows.append([str(v) for v in r.values()])
        elif isinstance(r, (list, tuple)):
            norm_rows.append(["" if v is None else str(v) for v in r])
        else:
            norm_rows.append([str(r)])

    # If preview columns provided use them; otherwise infer from data
    if columns:
        n_cols = max(1, len(columns))
    else:
        # infer columns as max row length
        n_cols = max((len(r) for r in norm_rows), default=1)

    # truncate data rows to fit on slide (reserve one row for header if columns present)
    max_data_rows = max_rows - (1 if columns else 0)
    if max_data_rows < 0:
        max_data_rows = 0
    display_rows = norm_rows[:max_data_rows]

    n_rows = len(display_rows) + (1 if columns else 0)
    if n_rows == 0:
        # nothing to show
        return

    # Create table: rows x cols
    # Use reasonable slide area
    left, top, width, height = Inches(0.5), Inches(1.2), Inches(8.5), Inches(3.5)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Fill header if present
    if columns:
        for ci in range(n_cols):
            text = columns[ci] if ci < len(columns) else ""
            table.cell(0, ci).text = str(text)

    # Fill body safely (guard indices)
    start_row = 1 if columns else 0
    for ri, row in enumerate(display_rows):
        for ci in range(n_cols):
            text = row[ci] if ci < len(row) else ""
            table.cell(start_row + ri, ci).text = str(text)

    return table


#Exporting data to Powerpoint
def generate_ppt(entry) -> BytesIO:
    """
    Generate a PowerPoint for a session entry which contains `messages`:
    entry["messages"] = [{"role":"turn","user_prompt":..., "assistant_run": {...}, "timestamp":...}, ...]
    Returns BytesIO.
    """
    prs = Presentation()
    layout = prs.slide_layouts[5]

    # Title slide for the session
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Agentic AI Report"
    session_title = entry.get("title") or (entry.get("prompt") or "")
    slide.placeholders[1].text = f"Session: {session_title}"
    created = entry.get("created_at") or entry.get("timestamp") or ""
    if created:
        # add a small subtitle for created time if available
        try:
            subtitle = slide.placeholders[1]
            subtitle.text += f"\nCreated: {created}"
        except Exception:
            pass

    # If messages absent (defensive) - fallback to single-run fields (but user said messages always present)
    messages = entry.get("messages", [])
    if not messages:
        # create a synthetic single-turn message using top-level entry fields
        messages = [{
            "role": "turn",
            "user_prompt": entry.get("prompt", ""),
            "assistant_run": {
                "prompt": entry.get("prompt"),
                "route": entry.get("route"),
                "result": entry.get("result"),
                "sql_query": entry.get("sql_query"),
                "preview_df": entry.get("preview_df"),
                "preview_df_columns": entry.get("preview_df_columns"),
                "header_candidate": entry.get("header_candidate"),
                "table_candidate_index": entry.get("table_candidate_index"),
                "header_updated": entry.get("header_updated"),
                "table_index_updated": entry.get("table_index_updated"),
                "updated_doc_path": entry.get("updated_doc_path"),
                "updated_doc_key": entry.get("updated_doc_key"),
                "web_links": entry.get("web_links"),
                "general_summary": entry.get("general_summary"),
                "comparison_summary": entry.get("comparison_summary"),
                "chart_info": entry.get("chart_info"),
                "faiss_summary": entry.get("faiss_summary"),
                "faiss_sources": entry.get("faiss_sources"),
                "faiss_images": entry.get("faiss_images"),
            },
            "timestamp": entry.get("timestamp")
        }]

    # Iterate through turns in stored order (do not change order)
    for idx, turn in enumerate(messages, start=1):
        user_prompt = turn.get("user_prompt") or ""
        timestamp = turn.get("timestamp") or ""
        assistant_run = turn.get("assistant_run") or {}

        # 1) Slide for the user prompt
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Turn {idx}: User Prompt"
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = user_prompt
        p.font.size = Pt(13)
        if timestamp:
            p = tf.add_paragraph()
            p.text = f"‚è± {timestamp}"
            p.font.size = Pt(10)

        # If assistant_run is empty, skip assistant slides
        if not assistant_run:
            continue

        route = assistant_run.get("route")

        # --- Document related (document route) ---
        if route == "document":
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"Turn {idx}: Document Update Summary"

            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
            tf = box.text_frame
            tf.word_wrap = True

            header = assistant_run.get("header_updated") or assistant_run.get("header_candidate")
            table_idx = assistant_run.get("table_index_updated") or assistant_run.get("table_candidate_index")
            updated_doc_path = assistant_run.get("updated_doc_path")

            p = tf.add_paragraph()
            p.text = f"Header Updated: {header or 'N/A'}"
            p.font.size = Pt(14)

            p = tf.add_paragraph()
            p.text = f"Table Index Updated: {table_idx if table_idx is not None else 'N/A'}"
            p.font.size = Pt(14)

            if updated_doc_path:
                p = tf.add_paragraph()
                p.text = f"Updated Document Path: {updated_doc_path}"
                p.font.size = Pt(12)

            # preview table BEFORE update
            preview_like = assistant_run.get("preview_df")
            preview_cols = assistant_run.get("preview_df_columns")
            prev_cols, prev_rows = [], []
            if isinstance(preview_like, pd.DataFrame):
                prev_cols, prev_rows = _rows_cols_from_serialized(preview_like)
            elif isinstance(preview_like, list) and preview_like:
                if preview_cols:
                    prev_cols = list(preview_cols)
                    prev_rows = [[("" if r.get(c) is None else str(r.get(c))) for c in prev_cols] for r in preview_like]
                else:
                    prev_cols, prev_rows = _rows_cols_from_serialized(preview_like)
            elif isinstance(preview_like, dict):
                prev_cols, prev_rows = _rows_cols_from_serialized(preview_like)

            if prev_rows:
                _add_table_slide(prs, f"Turn {idx}: Matched Table (Preview)", prev_cols, prev_rows, max_rows=7)

        # --- SQL query slide ---
        if assistant_run.get("sql_query"):
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"Turn {idx}: SQL Query"
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = assistant_run.get("sql_query")
            p.font.size = Pt(11)

        # --- SQL Result (if any) ---
        result = assistant_run.get("result")
        df_result = None
        if isinstance(result, list):
            try:
                df_result = pd.DataFrame(result)
            except Exception:
                df_result = None
        elif isinstance(result, pd.DataFrame):
            df_result = result

        if df_result is not None and not df_result.empty and route in ["sql", "document", "comp"]:
            # Add a table slide for SQL results, cap rows
            cols, rows = _rows_cols_from_serialized(df_result)
            if rows:
                _add_table_slide(prs, f"Turn {idx}: SQL Results", cols, rows, max_rows=6)

        # --- Comparison / General summaries ---
        if assistant_run.get("comparison_summary"):
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"Turn {idx}: Comparison Summary"
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
            tf = box.text_frame
            tf.word_wrap = True
            for para in str(assistant_run.get("comparison_summary")).split("\n"):
                if para.strip():
                    p = tf.add_paragraph()
                    p.text = para.strip()
                    p.font.size = Pt(12)

        if assistant_run.get("general_summary"):
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"Turn {idx}: General Summary"
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
            tf = box.text_frame
            tf.word_wrap = True
            for para in str(assistant_run.get("general_summary")).split("\n"):
                if para.strip():
                    p = tf.add_paragraph()
                    p.text = para.strip()
                    p.font.size = Pt(12)

        # --- Reconciliation slides (Two-Excel flow) ---
        if assistant_run.get("reconcile_df") or assistant_run.get("variance_commentary") or assistant_run.get("variance_summary"):
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"Turn {idx}: Reconciliation Summary"

            # Textbox for provenance and basic metadata
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.6))
            tf = box.text_frame
            tf.word_wrap = True

            # Source files / timestamp
            srcs = assistant_run.get("reconcile_source_files") or assistant_run.get("reconcile_source_files", [])
            if srcs:
                p = tf.paragraphs[0]
                p.text = "Source files: " + ", ".join([os.path.basename(str(s)) for s in srcs])
                p.font.size = Pt(11)
            else:
                p = tf.paragraphs[0]
                p.text = f"Reconciliation run: {assistant_run.get('timestamp', '')}"
                p.font.size = Pt(11)

            # Variance commentary slide (if present) ‚Äî short narrative from LLM
            commentary = assistant_run.get("variance_commentary")
            if commentary:
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = f"Turn {idx}: Variance Commentary"
                box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(3.5))
                tf = box.text_frame
                tf.word_wrap = True
                for para in str(commentary).split("\n"):
                    if para.strip():
                        p = tf.add_paragraph()
                        p.text = para.strip()
                        p.font.size = Pt(12)

            # Variance summary (key metrics) ‚Äî render as compact JSON-like text
            summary = assistant_run.get("variance_summary")
            if summary:
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = f"Turn {idx}: Variance Summary"
                box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(2.2))
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = "Top summary metrics:"
                p.font.size = Pt(12)
                try:
                    # pretty-print top-level keys
                    summary_text = json.dumps(summary, indent=2)
                    p = tf.add_paragraph()
                    p.text = summary_text[:3000]  # clip to reasonable length
                    p.level = 1
                    p.font.size = Pt(10)
                except Exception:
                    p = tf.add_paragraph()
                    p.text = str(summary)[:3000]
                    p.font.size = Pt(10)

            # Main variance / reconciliation table: convert serialized forms as needed and add table slide
            recon = assistant_run.get("reconcile_df")
            if recon:
                try:
                    # If stored as list-of-dicts or dict, coerce to cols/rows using helper
                    if isinstance(recon, (list, dict)):
                        cols, rows = _rows_cols_from_serialized(recon)
                        if rows:
                            _add_table_slide(prs, f"Turn {idx}: Variance Table (top rows)", cols, rows, max_rows=6)
                    elif isinstance(recon, pd.DataFrame):
                        cols, rows = _rows_cols_from_serialized(recon)
                        if rows:
                            _add_table_slide(prs, f"Turn {idx}: Variance Table (top rows)", cols, rows, max_rows=6)
                    else:
                        # try best-effort coercion
                        cols, rows = _rows_cols_from_serialized(recon)
                        if rows:
                            _add_table_slide(prs, f"Turn {idx}: Variance Table (top rows)", cols, rows, max_rows=6)
                except Exception:
                    # safe fallback: include a short note if table could not be rendered
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = f"Turn {idx}: Variance Table"
                    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
                    tf = box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = "Variance table present but could not be rendered in PPT. Please download the CSV from the session artifacts."
                    p.font.size = Pt(11)
            
        # --- Missing combinations slides (if present) ---
        missing_a = assistant_run.get("missing_in_A")
        missing_b = assistant_run.get("missing_in_B")

        # Missing in A (present in B)
        if missing_a:
            try:
                cols_a, rows_a = _rows_cols_from_serialized(missing_a)
                if rows_a:
                    _add_table_slide(prs, f"Turn {idx}: Missing in A (present in B) - top rows", cols_a, rows_a, max_rows=8)
                else:
                    # If no rows after serialization, add a note slide
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = f"Turn {idx}: Missing in A (present in B)"
                    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.0))
                    tf = box.text_frame
                    tf.text = "No example rows available for Missing-in-A."
            except Exception:
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = f"Turn {idx}: Missing in A (present in B)"
                box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
                tf = box.text_frame
                tf.word_wrap = True
                tf.text = "Missing-in-A table present but could not be rendered. Please download the CSV from the session artifacts."

        # Missing in B (present in A)
        if missing_b:
            try:
                cols_b, rows_b = _rows_cols_from_serialized(missing_b)
                if rows_b:
                    _add_table_slide(prs, f"Turn {idx}: Missing in B (present in A) - top rows", cols_b, rows_b, max_rows=8)
                else:
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = f"Turn {idx}: Missing in B (present in A)"
                    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.0))
                    tf = box.text_frame
                    tf.text = "No example rows available for Missing-in-B."
            except Exception:
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = f"Turn {idx}: Missing in B (present in A)"
                box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
                tf = box.text_frame
                tf.word_wrap = True
                tf.text = "Missing-in-B table present but could not be rendered. Please download the CSV from the session artifacts."


        # --- Web links (search/comp) ---
        web_links = assistant_run.get("web_links") or assistant_run.get("result") if route == "search" else assistant_run.get("web_links")
        if web_links:
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"Turn {idx}: Top Web Links"
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
            tf = box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(web_links, 1):
                # item could be tuple (markdown_link, summary) or simple string
                link_md, summary = (item[0], item[1]) if (isinstance(item, (list, tuple)) and len(item) >= 2) else (str(item), "")
                match = re.match(r"\[(.*?)\]\((.*?)\)", str(link_md))
                if match:
                    title, url = match.groups()
                else:
                    title, url = f"Link {i}", str(link_md)

                p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"{i}. {title}"
                try:
                    run.font.size = Pt(12)
                    run.hyperlink.address = url
                except Exception:
                    pass
                if summary:
                    s_p = tf.add_paragraph()
                    s_p.text = f"    ‚Ü≥ {str(summary)[:300]}"
                    s_p.font.size = Pt(11)

        # --- FAISS route slides if present (assistant_run or entry-level) ---
        faiss_summary = assistant_run.get("faiss_summary")
        faiss_sources = assistant_run.get("faiss_sources") or assistant_run.get("faiss_sources", [])
        faiss_images = assistant_run.get("faiss_images") or assistant_run.get("faiss_images", [])

        if faiss_summary:
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"Turn {idx}: FAISS Summary"
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
            tf = box.text_frame
            tf.word_wrap = True
            for para in str(faiss_summary).split("\n"):
                if para.strip():
                    p = tf.add_paragraph()
                    p.text = para.strip()
                    p.font.size = Pt(12)

        if faiss_sources:
            for i, src in enumerate(faiss_sources, 1):
                try:
                    docname, snippet, path = src[0], src[1], src[2] if len(src) >= 3 else (src[0], src[1], None)
                except Exception:
                    docname, snippet, path = str(src), "", None
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = f"Turn {idx}: FAISS Source {i} - {os.path.basename(path) if path else docname}"
                box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
                tf = box.text_frame
                tf.word_wrap = True
                for para in str(snippet).split("\n"):
                    if para.strip():
                        p = tf.add_paragraph()
                        p.text = para.strip()
                        p.font.size = Pt(11)

        if faiss_images and faiss_sources:
            # Only include images from the most-similar doc (first in faiss_sources)
            top_docname = faiss_sources[0][0] if isinstance(faiss_sources[0], (list, tuple)) else faiss_sources[0]
            top_images = [img for img in faiss_images if img.get("original_doc") == top_docname]
            if top_images:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = f"Turn {idx}: Images from {top_docname}"
                left = Inches(0.8)
                top = Inches(2.2)
                image_width = Inches(5.5)
                spacing = Inches(0.5)
                for im_meta in top_images:
                    img_path = im_meta.get("extracted_image_path")
                    if img_path and os.path.exists(img_path):
                        slide.shapes.add_picture(img_path, left, top, width=image_width)
                        top += Inches(3.2)
                        if top > Inches(6.5):
                            top = Inches(2.2)
                            left += image_width + spacing

        # --- Charts: if there is chart_info (you can expand how to render charts later) ---
        chart_info = assistant_run.get("chart_info")
        if chart_info:
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"Turn {idx}: Chart Info"
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
            tf = box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = str(chart_info)[:1500]

    # End: return PPT as BytesIO
    ppt_bytes = BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes


def _get_entry_datetime(entry):
    """
    Return a datetime for a history `entry`:
    - checks keys in order: 'timestamp', 'created_at', 'archived_at'
    - if not found, tries first/last message timestamps in entry['messages']
    - if still not found, returns current datetime
    Accepts string timestamps in format "%d %b %Y, %I:%M %p" or ISO format.
    """
    # 1) top-level fields
    ts = entry.get("timestamp") or entry.get("created_at") or entry.get("archived_at")
    # 2) try messages list
    if not ts:
        msgs = entry.get("messages") or []
        if msgs:
            # prefer first message timestamp then last
            ts = msgs[0].get("timestamp") or msgs[-1].get("timestamp") or msgs[0].get("assistant_run", {}).get("timestamp")
    # 3) fallback to now
    if not ts:
        return datetime.now()

    # 4) parse
    if isinstance(ts, datetime):
        return ts
    if isinstance(ts, (int, float)):
        try:
            return datetime.fromtimestamp(ts)
        except Exception:
            return datetime.now()

    ts_str = str(ts)
    # try your expected format first
    try:
        return datetime.strptime(ts_str, "%d %b %Y, %I:%M %p")
    except Exception:
        pass
    # try ISO formats
    try:
        return datetime.fromisoformat(ts_str)
    except Exception:
        pass
    # try common alternative formats (best-effort)
    alt_formats = [
        "%Y-%m-%d %H:%M:%S",
        "%Y-%m-%dT%H:%M:%S",
        "%d-%m-%Y %H:%M:%S",
    ]
    for fmt in alt_formats:
        try:
            return datetime.strptime(ts_str, fmt)
        except Exception:
            continue
    # give up -> return now
    return datetime.now()


def _format_dataframe_for_display(result_obj):
    """Helper: convert serialized list/dict to DataFrame and format numeric columns."""
    df = result_obj
    if isinstance(result_obj, list):
        df = pd.DataFrame(result_obj)
    if isinstance(df, pd.DataFrame):
        formatted_df = df.copy()
        for col in formatted_df.select_dtypes(include='number').columns:
            col_lower = col.lower()
            if "ratio" in col_lower:
                formatted_df[col] = formatted_df[col].apply(lambda x: f"{x * 100:.2f}%" if pd.notnull(x) else "")
            elif any(keyword in col_lower for keyword in money_keywords):
                formatted_df[col] = formatted_df[col].apply(lambda x: f"{x:,.0f}" if pd.notnull(x) else "")
        return formatted_df
    return result_obj

def _render_document_block(run):
    """Render document-related parts for a run_record (document route)."""
    header = run.get("header_updated") or run.get("header_candidate")
    table_idx = run.get("table_index_updated") or run.get("table_candidate_index")
    if header:
        st.markdown(f"**Header:** {header}")
    if table_idx is not None:
        st.markdown(f"**Table index:** {table_idx}")

    # preview (before update)
    preview_rows = run.get("preview_df")
    preview_cols = run.get("preview_df_columns")
    if preview_rows:
        st.markdown("**Matched Table (preview)**")
        try:
            preview_df = pd.DataFrame(preview_rows, columns=preview_cols if preview_cols else None)
        except Exception:
            preview_df = pd.DataFrame(preview_rows)
        st.dataframe(preview_df)

    # Download button - use stored unique key
    updated_path = run.get("updated_doc_path")
    download_key = run.get("updated_doc_key") or f"history_download_{hash(str(updated_path))}"
    if updated_path and os.path.exists(updated_path):
        with open(updated_path, "rb") as f:
            st.download_button(
                "üì• Download Updated Document",
                f,
                file_name=os.path.basename(updated_path) or "updated.docx",
                key=download_key
            )
    else:
        st.warning("Updated document file not available (may be temporary).")

def _render_faiss_block(entry):
    """Render faiss route elements from the archived entry (keeps original logic)."""
    st.subheader("üìò Internal Knowledge Base Answer:")
    st.markdown(entry.get("faiss_summary", "_No summary available._"))

    faiss_images = entry.get("faiss_images", [])
    faiss_sources = entry.get("faiss_sources", [])
    if faiss_images and faiss_sources:
        top_doc = faiss_sources[0][0]
        st.subheader(f"üñºÔ∏è Images from: {top_doc}")
        for meta in faiss_images:
            if meta.get("original_doc") == top_doc:
                img_path = meta.get("extracted_image_path")
                if img_path and os.path.exists(img_path):
                    st.image(img_path, caption=meta.get("caption", ""), use_container_width=True)

    st.subheader("üìÑ Document Sources:")
    base_dir = os.path.dirname(__file__)
    for i, (docname, snippet, path) in enumerate(faiss_sources, 1):
        col1, col2 = st.columns([0.85, 0.15])
        with col1:
            st.markdown(f"**{i}. {docname}**\n\n{snippet}")
        with col2:
            if path:
                full_path = os.path.join(base_dir, path).replace("\\", "/")
                if os.path.exists(full_path):
                    with open(full_path, "rb") as f:
                        st.download_button(
                            label="‚¨áÔ∏è",
                            data=f,
                            file_name=os.path.basename(path),
                            key=f"download_history_{i}"
                        )

def _render_run_by_route(run):
    import pandas as pd
    """Render the assistant_run (the full per-turn run_record) according to its route.
       This mirrors the previous single-entry rendering but scoped to a run record."""
    route = run.get("route")
    # For SQL-like or document-like routes, we show sql and result
    if route in ["sql", "document", "comp"]:
        # Document-specific extra (only for 'document' route)
        if route == "document":
            st.subheader("üìÑ Document Update (turn)")
            _render_document_block(run)

        # SQL Query Result
        if run.get("sql_query"):
            st.subheader("üßæ SQL Query:")
            st.code(run["sql_query"], language="sql")

            st.subheader("SQL Query Result:")
            result_df = run.get("result")
            formatted = _format_dataframe_for_display(result_df)
            if isinstance(formatted, pd.DataFrame):
                st.dataframe(formatted)
            else:
                st.text(formatted if formatted is not None else "_No result returned_")

        # For comparison runs, show summaries and web links if present
        if route == "comp":
             # If reconciliation outputs exist (two-excel flow), show those first
            comp_type = run.get("comp_type", "SQL-WEB")   
            #if run.get("reconcile_df") or run.get("reconcile_source_files"):
            if comp_type == "EXCEL-EXCEL":
                st.subheader("üîÅ Reconciliation (Two Excel Files)")

                # Show provenance / source files if available
                # srcs = run.get("reconcile_source_files") or []
                # if srcs:
                #     st.markdown("**Source files:**")
                #     for i, s in enumerate(srcs, 1):
                #         st.markdown(f"- {i}. `{s}`")

                # Variance commentary (LLM-generated) ‚Äî show first for quick context

                # Missing combinations in A / B
                missing_a = run.get("missing_in_A")
                missing_b = run.get("missing_in_B")
                if missing_a:
                    st.subheader("‚ö†Ô∏è Missing combinations in A (present in B)")
                    formatted_a = _format_dataframe_for_display(missing_a)
                    if isinstance(formatted_a, pd.DataFrame):
                        st.dataframe(formatted_a)
                        try:
                            csv_a = formatted_a.to_csv(index=False).encode("utf-8")
                            st.download_button("Download missing_in_A (CSV)", data=csv_a, file_name="missing_in_A.csv", mime="text/csv")
                        except Exception:
                            pass
                    else:
                        st.write(formatted_a)
                if missing_b:
                    st.subheader("‚ö†Ô∏è Missing combinations in B (present in A)")
                    formatted_b = _format_dataframe_for_display(missing_b)
                    if isinstance(formatted_b, pd.DataFrame):
                        st.dataframe(formatted_b)
                        try:
                            csv_b = formatted_b.to_csv(index=False).encode("utf-8")
                            st.download_button("Download missing_in_B (CSV)", data=csv_b, file_name="missing_in_B.csv", mime="text/csv")
                        except Exception:
                            pass
                    else:
                        st.write(formatted_b)


                # Main variance / reconciliation table
                recon_df = run.get("reconcile_df")
                if recon_df is not None:
                    st.subheader("üìã Variance Table (by combination)")
                    # formatted may already be serialized (list-of-dicts) in history; convert back if needed
                    try:
                        import pandas as pd
                        if isinstance(recon_df, list):
                            display_df = pd.DataFrame(recon_df)
                        elif isinstance(recon_df, dict):
                            display_df = pd.DataFrame([recon_df])
                        elif isinstance(recon_df, pd.DataFrame):
                            display_df = recon_df
                        else:
                            display_df = _format_dataframe_for_display(recon_df)
                    except Exception:
                        display_df = _format_dataframe_for_display(recon_df)

                    if isinstance(display_df, pd.DataFrame):
                        # üîπ Keep only rows with at least one non-zero numeric variance
                        numeric_cols = display_df.select_dtypes(include="number").columns
                        if len(numeric_cols) > 0:
                            # Replace NaN with 0 for the check, then keep any row where a numeric col != 0
                            filtered_df = display_df.loc[(display_df[numeric_cols].fillna(0) != 0).any(axis=1)]
                        else:
                            filtered_df = display_df  # no numeric columns, show as-is

                        st.dataframe(filtered_df)
                        # Provide CSV download for the user
                        try:
                            csv_bytes = filtered_df.to_csv(index=False).encode("utf-8")
                            st.download_button("Download Variance Table (CSV)", data=csv_bytes, file_name="reconciliation_variance.csv", mime="text/csv")
                        except Exception:
                            pass
                    else:
                        st.write(display_df)


                # Variance summary
                summary = run.get("variance_summary")
                if summary:
                    st.subheader("üìä Variance Summary")
                    try:
                        st.json(summary)
                    except Exception:
                        st.write(summary)


                # Variance Commentary
                commentary = run.get("variance_commentary")
                if commentary:
                    st.subheader("üìù Variance Commentary")
                    st.markdown(commentary)

                
            #SQL vs WEB
            else:
                if run.get("general_summary"):
                    st.subheader("üß† General Summary:")
                    st.markdown(run["general_summary"])
                st.subheader("üîó Top Web Links:")
                for i, (link, summary) in enumerate(run.get("web_links") or [], 1):
                    st.markdown(f"**{i}.** {link}")
                    st.markdown(f"_Summary:_\n{summary}")
                if run.get("comparison_summary"):
                    st.subheader("üÜö Comparison Summary:")
                    st.markdown(run["comparison_summary"])

    elif route == "faissdb":
        # faissdb runs may be stored within a run or (for older entries) at the top-level entry.
        _render_faiss_block(run)

    elif route == "search":
        if run.get("general_summary"):
            st.subheader("üß† General Summary:")
            st.markdown(run["general_summary"])
        st.subheader("üîó Top Web Links:")
        for i, (link, summary) in enumerate(run.get("result") or [], 1):
            st.markdown(f"**{i}.** {link}")
            st.markdown(f"_Summary:_\n{summary}")

    else:
        # Fallback: print any summary or raw result
        if run.get("general_summary"):
            st.subheader("üß† General Summary:")
            st.markdown(run["general_summary"])
        if run.get("result"):
            st.subheader("Result:")
            formatted = _format_dataframe_for_display(run.get("result"))
            if isinstance(formatted, pd.DataFrame):
                st.dataframe(formatted)
            else:
                st.text(formatted)


# ---- LangGraph Setup ----
graph_builder = StateGraph(GraphState)
graph_builder.add_node("router", RouterNode())
graph_builder.add_node("vanna_sql", vanna_node)
graph_builder.add_node("serp_search", serp_node)
graph_builder.add_node("doc_update", document_node)
graph_builder.add_node("comp", comp_node)
graph_builder.add_node("faissdb", faissdb_node)

def router_logic(state: GraphState):
    if state['route'] == 'sql': return "vanna_sql"
    elif state['route'] == 'search': return "serp_search"
    elif state['route'] == 'document': return "doc_update"
    elif state['route'] == 'comp': return "comp"
    elif state['route'] == 'faissdb': return "faissdb"
    else: return END    

graph_builder.set_entry_point("router")

# ‚úÖ Execution routing
graph_builder.add_conditional_edges("router", router_logic)

# ‚úÖ Visualization support ‚Äî add all potential router paths
#graph_builder.add_edge("router", "vanna_sql")
#graph_builder.add_edge("router", "serp_search")
#graph_builder.add_edge("router", "doc_update")

# Regular path to end
graph_builder.add_edge("vanna_sql", END)
graph_builder.add_edge("serp_search", END)
graph_builder.add_edge("doc_update", END)
graph_builder.add_edge("comp", END)
graph_builder.add_edge("faissdb", END)

agent_graph = graph_builder.compile()

# ---- Streamlit UI ----
st.title("\U0001F9E0 Agentic AI Assistant (Insurance)")


def format_date_label(chat_date: date) -> str:
    today = date.today()
    if chat_date == today:
        return "Today"
    elif chat_date == today - timedelta(days=1):
        return "Yesterday"
    else:
        return chat_date.strftime("%d %b %Y")
    
def generate_title(prompt: str) -> str:
    try:
        title_prompt = f"Summarize the following user query into a short title:\n\n'{prompt}'\n\nKeep it under 7 words."
        return call_llm(title_prompt)
    except:
        return prompt[:40] + ("..." if len(prompt) > 40 else "")
    

# ‚úÖ Initialize chat history and active index
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "active_chat_index" not in st.session_state:
    st.session_state.active_chat_index = None

# NEW: the live session (holds a list of Q/A messages)
if "current_session" not in st.session_state:
    st.session_state.current_session = {
        "id": str(uuid.uuid4()),
        "title": None,
        "created_at": datetime.now().strftime("%d %b %Y, %I:%M %p"),
        "messages": []  # each message: {"role":"user"/"assistant","text":..., "route":..., "result": ...}
    }

# ‚úÖ Sidebar: Clear + View + Export
with st.sidebar:
    st.header("üóÇÔ∏è Session")
    if st.button("üßπ Clear Chat History"):
        st.session_state.chat_history = []
        st.session_state.active_chat_index = None
        st.success("Chat history cleared!")

# ‚úÖ Group and render chat history in sidebar
grouped = {}
for chat in st.session_state.chat_history:
    chat_dt = _get_entry_datetime(chat)
    chat_date = chat_dt.date()
    grouped.setdefault(chat_date, []).append(chat)

for group_date in sorted(grouped.keys(), reverse=True):
    label = format_date_label(group_date)
    with st.sidebar.expander(f"üìÖ {label}"):
        entries = sorted(grouped[group_date], key=lambda e: not e.get("pinned", False))
        for idx, chat in enumerate(entries):
            title = chat.get("title") or chat["prompt"][:40]
            pin_icon = "üìå " if chat.get("pinned") else ""
            if st.button(f"{pin_icon}{title}", key=f"chat_{group_date}_{idx}"):
                st.session_state.active_chat_index = st.session_state.chat_history.index(chat)
                st.session_state.user_prompt = chat["prompt"]
                st.session_state.just_ran_agent = False

# ‚úÖ Export chat history

def safe_serialize_obj(obj):
    """
    Convert obj to a JSON-serializable representation.
    Handles: pandas.DataFrame, list[dict], numpy types, datetime/date, Path, BytesIO.
    For unknown objects, falls back to str(obj).
    """
    # None / primitives
    if obj is None or isinstance(obj, (str, bool, int, float)):
        return obj

    # datetime / date
    if isinstance(obj, (datetime, date)):
        return obj.isoformat()

    # Path
    if isinstance(obj, Path):
        return str(obj)

    # BytesIO -> base64 (optional; careful about size)
    if isinstance(obj, BytesIO):
        # convert to base64 string (avoid if large)
        obj.seek(0)
        b64 = base64.b64encode(obj.read()).decode("ascii")
        return {"__bytes_base64__": b64}

    # pandas DataFrame
    if isinstance(obj, pd.DataFrame):
        try:
            # prefer records orient (list of dicts)
            return obj.where(pd.notnull(obj), None).to_dict(orient="records")
        except Exception:
            # fallback to string
            return str(obj)

    # pandas Series
    if isinstance(obj, pd.Series):
        try:
            return obj.where(pd.notnull(obj), None).to_dict()
        except Exception:
            return list(obj)

    # numpy scalar types
    if isinstance(obj, (np.generic,)):
        return obj.item()

    # lists / tuples / sets
    if isinstance(obj, (list, tuple, set)):
        return [safe_serialize_obj(i) for i in obj]

    # dicts -> apply recursively
    if isinstance(obj, dict):
        out = {}
        for k, v in obj.items():
            # ensure keys are strings
            key = str(k)
            out[key] = safe_serialize_obj(v)
        return out

    # dataclasses? try to convert to dict
    # fallback: try to use __dict__ if present
    if hasattr(obj, "__dict__"):
        try:
            return safe_serialize_obj(vars(obj))
        except Exception:
            pass

    # last resort: stringify
    try:
        return str(obj)
    except Exception:
        return None


def serialize_chat_history(history):
    """
    Given st.session_state.chat_history (list of dicts), produce a JSON string safely.
    Use this instead of plain json.dumps(history).
    """
    safe_history = []
    for entry in history:
        safe_entry = {}
        # iterate keys in original entry and serialize values
        for k, v in entry.items():
            safe_entry[str(k)] = safe_serialize_obj(v)
        # ensure messages list (if present) is serialized as well
        if "messages" in safe_entry and isinstance(safe_entry["messages"], list):
            safe_messages = []
            for m in safe_entry["messages"]:
                safe_messages.append(safe_serialize_obj(m))
            safe_entry["messages"] = safe_messages
        safe_history.append(safe_entry)

    return json.dumps(safe_history, indent=2, ensure_ascii=False)

history_json = serialize_chat_history(st.session_state.chat_history)

st.download_button("‚¨áÔ∏è Export Chat History", history_json, file_name="chat_history.json")

# Render before running agent (all dashed)
#with st.expander("üß≠ Workflow Graph (Initial)"):
#    visualize_workflow(graph_builder)

# ‚úÖ Initialize just_ran_agent flag if not already
if "just_ran_agent" not in st.session_state:
    st.session_state.just_ran_agent = False

# ‚úÖ UI Control Logic: if user is NOT viewing past chat
if st.session_state.active_chat_index is None:


    # near the main input area (when active_chat_index is None i.e. live session)
    if st.button("Start New Session"):
        sess = st.session_state.current_session
        if not sess["messages"]:
            # nothing to archive
            st.warning("Current session is empty ‚Äî nothing to archive.")
        else:
            last_run = sess["messages"][-1]["assistant_run"]

            entry = {
                "id": sess["id"],
                "title": sess.get("title") or generate_title(sess["messages"][0]["user_prompt"]),
                "prompt": sess["messages"][0]["user_prompt"],   # first prompt in session
                "route": last_run.get("route"),
                "result": last_run.get("result"),
                "sql_query": last_run.get("sql_query"),

                # Document fields (from last run)
                "preview_df": last_run.get("preview_df"),
                "preview_df_columns": last_run.get("preview_df_columns"),
                "header_candidate": last_run.get("header_candidate"),
                "table_candidate_index": last_run.get("table_candidate_index"),
                "header_updated": last_run.get("header_updated"),
                "table_index_updated": last_run.get("table_index_updated"),
                "updated_doc_path": last_run.get("updated_doc_path"),
                "updated_doc_key": last_run.get("updated_doc_key"),

                # Web and summaries
                "web_links": last_run.get("web_links"),
                "general_summary": last_run.get("general_summary"),
                "comparison_summary": last_run.get("comparison_summary"),
                "chart_info": last_run.get("chart_info"),

                # FAISS stuff
                "faiss_summary": last_run.get("faiss_summary"),
                "faiss_sources": last_run.get("faiss_sources"),
                "faiss_images": last_run.get("faiss_images"),

                # keep the whole message list (full session)
                "messages": sess["messages"],

                # meta
                "created_at": sess.get("created_at"),
                "archived_at": datetime.now().strftime("%d %b %Y, %I:%M %p")
            }

            st.session_state.chat_history.append(entry)

            # Reset current_session to a fresh empty session
            st.session_state.current_session = {
                "id": str(uuid.uuid4()),
                "title": None,
                "created_at": datetime.now().strftime("%d %b %Y, %I:%M %p"),
                "messages": []
            }

            # Reset conversation context as well (since it's per-session)
            st.session_state.conversation = []
            st.session_state.active_chat_index = len(st.session_state.chat_history) - 1
            st.rerun()


    with st.form(key="query_form"):
        user_prompt = st.text_input("Enter your query:")
        #doc_file = st.file_uploader("Upload Insurance Document (.docx)", type=["docx"])
        uploaded_file1 = st.file_uploader("Attach first Excel file", type=["docx", "xlsx", "xls", "csv"], key="file1")
        uploaded_file2 = st.file_uploader("Attach second Excel file", type=["docx", "xlsx", "xls", "csv"], key="file2")
        submitted = st.form_submit_button("Run Agent")


    if submitted:
    # Only run when prompt is entered and changed
    #if user_prompt and (
    #    "last_prompt" not in st.session_state
    #    or st.session_state["last_prompt"] != user_prompt
    #):
        st.session_state["last_prompt"] = user_prompt

        # Persist uploaded files to temp paths so downstream nodes can use file paths
        file1_path = None
        file2_path = None
        uploaded_file1_is_excel = False
        uploaded_file2_is_excel = False

        if uploaded_file1 is not None:
            # Extract and normalize the file extension
            ext1 = uploaded_file1.name.split(".")[-1].lower().strip()
            
            with tempfile.NamedTemporaryFile(delete=False, suffix="." + (uploaded_file1.name.split(".")[-1])) as tmp1:
                tmp1.write(uploaded_file1.read())
                file1_path = tmp1.name
                
                # Set flag based on extension
                uploaded_file1_is_excel = ext1 in ["xls", "xlsx", "csv"]

        if uploaded_file2 is not None:
            ext2 = uploaded_file2.name.split(".")[-1].lower().strip()

            with tempfile.NamedTemporaryFile(delete=False, suffix="." + (uploaded_file2.name.split(".")[-1])) as tmp2:
                tmp2.write(uploaded_file2.read())
                file2_path = tmp2.name
                
                # Set flag based on extension
                uploaded_file2_is_excel = ext2 in ["xls", "xlsx", "csv"]

        # if doc_file:
        #     with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
        #         tmp.write(doc_file.read())
        #         doc_path = tmp.name
        # else:
        #     doc_path = None


        state: GraphState = {
            # Core inputs
            "user_prompt": user_prompt,
            "uploaded_file1": uploaded_file1,        # raw UploadedFile object (optional)
            "uploaded_file2": uploaded_file2,        # raw UploadedFile object (optional)
            "uploaded_file1_path": file1_path,       # filesystem path to temp file (or None)
            "uploaded_file2_path": file2_path,       # filesystem path to temp file (or None)
            "uploaded_file1_is_excel": uploaded_file1_is_excel,         # True if file1 appears excel/csv-like
            "uploaded_file2_is_excel": uploaded_file2_is_excel,         # True if file2 appears excel/csv-like
#            "doc_loaded": doc_path is not None,
#            "document_path": doc_path,

            # Prompts
            "vanna_prompt": None,
            "fuzzy_prompt": None,

            # Routing
            "route": None,

            # SQL-related
            "sql_result": None,
            "sql_query": None,

            # Document-related
            "updated_doc_path": None,
            "header_candidate": None,
            "table_candidate_index": None,
            "header_updated": None,
            "table_index_updated": None,
            "updated_doc_key": None,
            "preview_df": None,
            "preview_df_columns": None,

            #"candidate_tables": None,

            # External search
            "web_links": None,

            # Visualization / summaries
            "comparison_summary": None,
            "general_summary": None,
            "chart_info": None,

            # FAISS Knowledge base
            "faiss_summary": None,
            "faiss_sources": None,
            "faiss_images": None,

            # --- Reconciliation-specific placeholders ---
            # These will be filled by comp_node / reconciliation_node when two files are attached

            "reconcile_df": None,              # full variance table (pandas DataFrame)
            "missing_in_A": None,              # DataFrame or list-of-dicts
            "missing_in_B": None,              # DataFrame or list-of-dicts
            "variance_summary": None,          # dict with totals/aggregates
            "variance_commentary": None,       # LLM-generated narrative text
            "reconcile_source_files": [file1_path, file2_path],  # for provenance
            "comp_type": None
        }

        with st.spinner("Running Agent..."):
            try:
                output = agent_graph.invoke(state)
                st.session_state.output = output
                followups = generate_follow_up_questions(user_prompt)
                st.session_state.followups = followups

            except Exception as e:
                st.error(f"Agent crashed due to error: {e}")
                import traceback
                st.code(traceback.format_exc())
                st.stop()


        # Build a full "run_record" that mirrors your old entry exactly.
        # Note: keep this small enough to be serializable. Convert any DataFrames to list-of-dicts,
        # and for large files use file paths/keys rather than embedding binary content.
        def safe_serialize_preview_df(df_like):
            # If it's already a list-of-dicts, return as-is.
            # If it's a pandas DataFrame, convert to records.
            try:
                import pandas as pd
                if isinstance(df_like, pd.DataFrame):
                    return df_like.to_dict(orient="records")
            except Exception:
                pass
            return df_like

        run_record = {
            "id": str(uuid.uuid4()),
            "prompt": user_prompt,
            "title": generate_title(user_prompt),
            "route": output.get("route"),

            # Results
            "result": output.get("sql_result") if output.get("route") in ["sql", "document", "comp"] else output.get("web_links"),
            "sql_query": output.get("sql_query"),

            # Document-related fields
            "preview_df": safe_serialize_preview_df(output.get("preview_df")),
            "preview_df_columns": output.get("preview_df_columns"),
            "header_candidate": output.get("header_candidate"),
            "table_candidate_index": output.get("table_candidate_index"),
            "header_updated": output.get("header_updated"),
            "table_index_updated": output.get("table_index_updated"),
            "updated_doc_path": output.get("updated_doc_path"),
            "updated_doc_key": output.get("updated_doc_key"),
            #"candidate_tables": output.get("candidate_tables"),

            # External search
            "web_links": output.get("web_links"),

            # Summaries / visualization
            "general_summary": output.get("general_summary"),
            "comparison_summary": output.get("comparison_summary"),
            "chart_info": output.get("chart_info"),

            # FAISS Knowledge base
            "faiss_summary": output.get("faiss_summary"),
            "faiss_sources": output.get("faiss_sources"),
            "faiss_images": output.get("faiss_images"),

            # Reconciliation-specific serialized outputs (if available)
            "comp_type": output.get("comp_type"),
            "reconcile_df": safe_serialize_preview_df(output.get("reconcile_df")),
            "missing_in_A": safe_serialize_preview_df(output.get("missing_in_A")),
            "missing_in_B": safe_serialize_preview_df(output.get("missing_in_B")),
            "variance_summary": output.get("variance_summary"),
            "variance_commentary": output.get("variance_commentary"),
            "reconcile_source_files": output.get("reconcile_source_files", [file1_path, file2_path]),

            # any additional custom fields the nodes may emit
            "extra": output.get("extra", {}),

            # Meta
            "timestamp": datetime.now().strftime("%d %b %Y, %I:%M %p")
        }

        # Append the run_record to the live session's messages
        # We'll store both the user input and the assistant/run result as a single message item
        st.session_state.current_session["messages"].append({
            "role": "turn",                # 'turn' groups the user->assistant result; you can also use separate user/assistant items
            "user_prompt": user_prompt,
            "assistant_run": run_record,
            "timestamp": run_record["timestamp"]
        })

        # Mark that agent ran (for UI)
        st.session_state.just_ran_agent = True

        #Render workflow + live session (reverse-chronological, assistant-first) ----------
        col_left, col_mid, col_right = st.columns([4, 0.4, 1.5])

        with col_right:
            if st.session_state.get("output"):
                st.markdown("### üß≠ Workflow Diagram")
                visualize_workflow(graph_builder, active_route=st.session_state["output"].get("route"))

        with col_left:
            # Render the entire current session in reverse-chronological order (latest first).
            sess = st.session_state.get("current_session", {"messages": [], "created_at": ""})
            messages = sess.get("messages", [])

            st.markdown("### üí¨ Current Session (Live) ‚Äî Latest first")
            sess_title = sess.get("title") or (messages[0].get("user_prompt") if messages else "New Session")
            st.caption(f"Session: {sess_title} ‚Äî Created: {sess.get('created_at', '')}")

            if not messages:
                st.info("No messages yet. Ask your first question!")
            else:
                # Show latest turn first (assistant output first, then user prompt)
                rev_messages = list(reversed(messages))

                for display_idx, turn in enumerate(rev_messages, start=1):
                    assistant_run = turn.get("assistant_run")
                    assistant_text = None
                    if not assistant_run:
                        assistant_text = turn.get("assistant_text") or turn.get("result") or turn.get("answer")

                    user_prompt_cap = (turn.get("user_prompt") or "").strip()
                    # Assistant output first (most recent on top)
                    st.markdown(f"#### {display_idx}. User Prompt: {user_prompt_cap}")
                    if assistant_run:
                        # Preferred: reuse your robust renderer
                        try:
                            _render_run_by_route(assistant_run)

                            # get the user prompt for this turn
                            user_prompt = (turn.get("user_prompt") or "").strip()
                            # Only proceed if user prompt mentions plotting intent
                            plotting_keywords = ["plot", "draw", "visualize", "chart", "bar graph", "line graph", "pie chart", "graph"]
                            if user_prompt and any(word in user_prompt.lower() for word in plotting_keywords):
                                # Reconstruct a DataFrame from assistant_run result if possible
                                sql_df = None
                                res = assistant_run.get("result") if assistant_run else None

                                # result may be list-of-dicts or a DataFrame
                                if isinstance(res, list):
                                    try:
                                        sql_df = pd.DataFrame(res)
                                    except Exception:
                                        sql_df = None
                                elif isinstance(res, pd.DataFrame):
                                    sql_df = res
                                else:
                                    # sometimes result stored as stringified table or under a different key
                                    # try assistant_run.get("preview_df") as fallback (document route)
                                    preview = assistant_run.get("preview_df") if assistant_run else None
                                    if isinstance(preview, list):
                                        try:
                                            sql_df = pd.DataFrame(preview)
                                        except Exception:
                                            sql_df = None

                                # Only try to suggest/plot if we have a reasonable DataFrame
                                if sql_df is not None and not sql_df.empty:
                                    user_chart_type = None
                                    try:
                                        user_chart_type = get_user_chart_type(user_prompt)
                                    except Exception:
                                        # failing to parse user chart type shouldn't block plotting
                                        user_chart_type = None

                                    chart_info = None
                                    try:
                                        chart_info = suggest_chart(sql_df)
                                    except Exception as e:
                                        chart_info = None
                                        # optionally show debug: st.info(f"chart suggestion error: {e}")

                                    if chart_info and user_chart_type:
                                        chart_info["type"] = user_chart_type

                                    if chart_info:
                                        try:
                                            plot_chart(sql_df, chart_info)
                                        except Exception as e:
                                            st.warning(f"Could not render chart: {e}")
                                else:
                                    # no tabular data to chart
                                    st.info("No tabular result available in this turn to create a chart from.")

                        except Exception as e:
                            # fallback rendering if helper missing or errors out
                            st.warning(f"(Renderer error: {e}) Showing raw assistant summary/result instead.")
                            if assistant_run.get("general_summary"):
                                st.markdown(assistant_run.get("general_summary"))
                            elif assistant_run.get("result"):
                                st.write(assistant_run.get("result"))
                            elif assistant_run.get("web_links"):
                                for i, item in enumerate(assistant_run.get("web_links"), 1):
                                    st.markdown(f"{i}. {item}")
                    else:
                        if assistant_text:
                            st.markdown(assistant_text)
                        else:
                            st.text("_No assistant output available for this turn_")

                    # Then show corresponding user prompt below the assistant output
                    user_prompt = turn.get("user_prompt") or turn.get("text") or turn.get("prompt") or ""
                    if user_prompt:
                        st.markdown(f"**You:** {user_prompt}")

                    # Timestamp (if present)
                    ts = turn.get("timestamp") or (assistant_run.get("timestamp") if assistant_run else None)
                    if ts:
                        st.caption(f"üïí {ts}")

                    st.markdown("---")

            # After rendering the live conversation, show followups and the single PPT export button
            if st.session_state.get("followups"):
                st.markdown("### üí¨ You could also ask:")
                for q in st.session_state.get("followups", []):
                    st.markdown(f"- üëâ {q}")

            try:
                entry_for_export = {
                    "id": sess.get("id"),
                    "title": sess.get("title"),
                    "prompt": sess.get("messages")[0]["user_prompt"] if sess.get("messages") else sess.get("title") or "",
                    "messages": sess.get("messages", []),
                    "created_at": sess.get("created_at"),
                }

                # Generate PPT buffer and show download button directly (same as history)
                ppt_buffer = generate_ppt(entry_for_export)
                # ensure pointer at start
                try:
                    ppt_buffer.seek(0)
                except Exception:
                    pass

                st.download_button(
                    label="‚¨áÔ∏è Export current session to PPT",
                    data=ppt_buffer,
                    file_name="agentic_ai_session.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"download_ppt_live_{entry_for_export.get('id')}"
                )

            except Exception as e:
                st.warning(f"PPT export not available: {e}")

        # Keep the app state similar to before: mark that we've finished rendering post-run
        st.session_state.just_ran_agent = False
        st.session_state.active_chat_index = None


else:
    # ‚úÖ If user is viewing previous chat, show message + unlock option
    st.info("üìú You're viewing a previous conversation. Click below to start a new query.")
    if st.button("Start New Query"):
        st.session_state.active_chat_index = None
        st.session_state.user_prompt = ""

        st.rerun() 

# ‚úÖ History Rendering selected chat in main area
# ---------------------------
# Render selected chat in main area (full session support)
# ---------------------------

# ---------- Main rendering logic ----------
if st.session_state.active_chat_index is not None and not st.session_state.just_ran_agent:
    entry = st.session_state.chat_history[st.session_state.active_chat_index]

    # If the entry contains a 'messages' list (archived session), render each turn in order.
    if entry.get("messages"):
        # Session header: use title if present else first prompt
        title = entry.get("title") or entry.get("prompt") or entry["messages"][0].get("user_prompt", "Session")
        st.markdown(f"### üóÇÔ∏è Session: {title}")
        # show created/archived metadata if present
        created = entry.get("created_at") or entry.get("timestamp")
        archived = entry.get("archived_at")
        if created:
            st.caption(f"Created: {created}")
        if archived:
            st.caption(f"Archived: {archived}")

        # Iterate through messages in stored order (do not change order)
        for idx, turn in enumerate(entry["messages"], start=1):
            # Support two stored formats:
            # 1) new format: {"role":"turn","user_prompt":..., "assistant_run": run_record, "timestamp":...}
            # 2) older format: {"role":"user","text": "..."} or similar
            user_prompt = turn.get("user_prompt") or turn.get("text") or turn.get("prompt")
            timestamp = turn.get("timestamp") or turn.get("assistant_run", {}).get("timestamp") or ""
            st.markdown(f"**{idx}. You:** {user_prompt}")
            if timestamp:
                st.caption(f"üïí {timestamp}")

            # Fetch assistant run record (if present)
            assistant_run = turn.get("assistant_run")
            if assistant_run:
                # Render the assistant run preserving all fields and order
                _render_run_by_route(assistant_run)
            else:
                # Fallback: maybe the message stored assistant text in `turn['assistant_text']` or similar
                assistant_text = turn.get("assistant_text") or turn.get("answer") or turn.get("result")
                if assistant_text:
                    st.markdown(f"**Assistant:** {assistant_text}")

            # Divider between turns
            st.markdown("---")

        # After full session rendering, allow export to PPT (keeps previous behavior)
        try:
            ppt_buffer = generate_ppt(entry)
            st.download_button("‚¨áÔ∏è Export to PPT", ppt_buffer, file_name="agentic_ai_output.pptx")
        except Exception:
            # keep UI robust if PPT generator fails for some entries
            st.warning("Unable to export PPT for this session.")

    else:
        st.text("Message not found")
